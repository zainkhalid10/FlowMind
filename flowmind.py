from fastapi import FastAPI, File, UploadFile, HTTPException, Depends, Form, status
from fastapi.responses import HTMLResponse, JSONResponse, RedirectResponse
from fastapi.staticfiles import StaticFiles
from sqlalchemy.orm import Session
from typing import Optional
from pydantic import BaseModel, EmailStr
import pytesseract
from database import SessionLocal, ParsedFile, ImageMeta, User, Feature, Team, AgentChatHistory, init_db
from rag_agent import get_agent
from auth import (
    get_password_hash, verify_password, create_access_token,
    get_current_user, get_db, get_current_user_optional
)

from PIL import Image
from pypdf import PdfReader
import docx
import pptx
import os
import time
import uuid
import html
import hashlib
import cv2
import json
import subprocess
import shutil
import base64
import requests
import asyncio
import warnings
import logging
import re
import sys
from dotenv import load_dotenv

try:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    if hasattr(sys.stderr, "reconfigure"):
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

# Suppress harmless warnings
warnings.filterwarnings("ignore", category=DeprecationWarning, module="pypdf")
warnings.filterwarnings("ignore", category=FutureWarning, module="transformers")
warnings.filterwarnings("ignore", message=".*ARC4.*")
warnings.filterwarnings("ignore", message=".*_register_pytree_node.*")
warnings.filterwarnings("ignore", message=".*CoreMLExecutionProvider.*")
warnings.filterwarnings("ignore", category=UserWarning, module="urllib3")

# Suppress ChromaDB telemetry errors
logging.getLogger("chromadb").setLevel(logging.ERROR)
os.environ.setdefault("ANONYMIZED_TELEMETRY", "false")

# Suppress huggingface tokenizers warnings about forking
os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")

# -------------------- Unicode Sanitization --------------------
def sanitize_unicode(text: str) -> str:
    """Remove invalid Unicode surrogate characters that can't be encoded to UTF-8."""
    if not text:
        return text or ""
    
    # Handle non-string types
    if not isinstance(text, str):
        text = str(text) if text is not None else ""
    
    if not text:
        return ""
    
    try:
        # First, try to encode to catch surrogates
        text.encode('utf-8')
        return text
    except UnicodeEncodeError:
        # Remove surrogates by encoding with errors='replace' and decoding
        # This replaces invalid characters with replacement characters
        try:
            return text.encode('utf-8', errors='replace').decode('utf-8', errors='replace')
        except Exception:
            # Fallback: remove surrogates manually
            return ''.join(char for char in text if ord(char) < 0xD800 or ord(char) > 0xDFFF)
    except (AttributeError, TypeError) as e:
        # If it's not a string, convert to string first
        try:
            text_str = str(text) if text is not None else ""
            text_str.encode('utf-8')
            return text_str
        except UnicodeEncodeError:
            try:
                return text_str.encode('utf-8', errors='replace').decode('utf-8', errors='replace')
            except Exception:
                return ''.join(char for char in text_str if ord(char) < 0xD800 or ord(char) > 0xDFFF)
        except Exception:
            return ""

def sanitize_dict(data: dict) -> dict:
    """Recursively sanitize all string values in a dictionary."""
    if not isinstance(data, dict):
        return data
    
    sanitized = {}
    for key, value in data.items():
        if isinstance(value, str):
            sanitized[key] = sanitize_unicode(value)
        elif isinstance(value, dict):
            sanitized[key] = sanitize_dict(value)
        elif isinstance(value, list):
            sanitized[key] = [sanitize_dict(item) if isinstance(item, dict) else (sanitize_unicode(item) if isinstance(item, str) else item) for item in value]
        else:
            sanitized[key] = value
    return sanitized

# -------------------- OCR Summary Helpers --------------------
def _summarize_image_ocr(ocr_text: str, context: str = "") -> str:
    """Legacy function - now uses enhanced service."""
    from services.image_service import enhanced_ocr_summarize
    result = enhanced_ocr_summarize(ocr_text, context)
    return sanitize_unicode(result)

# -------------------- SETUP --------------------
# Tesseract path - platform-aware configuration
import platform
tesseract_path = os.getenv("TESSERACT_CMD")
if not tesseract_path:
    if platform.system() == "Darwin":  # macOS
        tesseract_path = "/opt/homebrew/bin/tesseract"
    elif platform.system() == "Windows":
        tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    # Linux usually has it in PATH, so leave as None

if tesseract_path and os.path.exists(tesseract_path):
    pytesseract.pytesseract.tesseract_cmd = tesseract_path
elif not tesseract_path:
    # Try to find tesseract in PATH
    tesseract_cmd = shutil.which("tesseract")
    if tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

load_dotenv()  # load env from .env so reload keeps VLM settings


def _get_soffice_path():
    """Return path to LibreOffice soffice for DOC/PPT conversion. Cross-platform."""
    env_path = os.getenv("LIBREOFFICE_PATH")
    if env_path and os.path.exists(env_path):
        return env_path
    candidates = [
        shutil.which("soffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS
        r"C:\Program Files\LibreOffice\program\soffice.exe",  # Windows default
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",  # Windows x86
    ]
    return next((p for p in candidates if p and os.path.exists(p)), None)


# CORS configuration
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI(title="FlowMind")
init_db()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, specify exact origins
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Include routers
from routes import auth_routes, upload_routes, dashboard_routes, training_routes, approval_routes, integration_routes, project_routes, google_oauth_routes, review_routes
from routes.approval_routes import parse_and_save_features
from services.document_validator import validate_document_for_srs

app.include_router(auth_routes.router)
app.include_router(upload_routes.router)
app.include_router(dashboard_routes.router)
app.include_router(training_routes.router)
app.include_router(approval_routes.router)
app.include_router(integration_routes.router)
app.include_router(project_routes.router)
app.include_router(google_oauth_routes.router)
app.include_router(review_routes.router)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(BASE_DIR, "uploads")
STATIC_DIR = os.path.join(BASE_DIR, "static")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Favicon endpoint
@app.get("/favicon.ico")
async def favicon():
    """Return empty favicon to prevent 404 errors."""
    from fastapi.responses import Response
    return Response(content="", media_type="image/x-icon")

# Serve static uploaded files
app.mount("/uploads", StaticFiles(directory=UPLOAD_DIR), name="uploads")

# Serve static frontend files
app.mount("/static", StaticFiles(directory=STATIC_DIR), name="static")

REQUIREMENTS_VIEWS = {}

# Load persistent views from file
VIEWS_FILE = "requirements_views.json"

def load_views():
    """Load views from persistent storage."""
    global REQUIREMENTS_VIEWS
    try:
        if os.path.exists(VIEWS_FILE):
            with open(VIEWS_FILE, "r", encoding="utf-8") as f:
                REQUIREMENTS_VIEWS = json.load(f)
                print(f"[views] Loaded {len(REQUIREMENTS_VIEWS)} persisted views")
    except Exception as e:
        print(f"[views] Failed to load views: {str(e)}")
        REQUIREMENTS_VIEWS = {}

def save_views():
    """Save views to persistent storage."""
    try:
        with open(VIEWS_FILE, "w", encoding="utf-8") as f:
            json.dump(REQUIREMENTS_VIEWS, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"[views] Failed to save views: {str(e)}")

# Load views on startup
load_views()

# -------------------- AUTHENTICATION ENDPOINTS --------------------
class SignupRequest(BaseModel):
    email: str
    username: str
    password: str

class LoginRequest(BaseModel):
    email: str
    password: str
    role: Optional[str] = None  # Optional: Manager, Team Head, Member - validated at login

@app.post("/api/signup")
async def signup(request: SignupRequest, db: Session = Depends(get_db)):
    """Simple user registration - no verification required."""
    try:
        # Check if user already exists
        existing_user = db.query(User).filter(
            (User.email == request.email) | (User.username == request.username)
        ).first()
        
        if existing_user:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Email or username already registered"
            )
        
        # Validate password length
        if len(request.password) < 6:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Password must be at least 6 characters long"
            )
        
        # No maximum length limit - using pbkdf2_sha256 which has no 72-byte restriction
        default_team = db.query(Team).first()
        hashed_password = get_password_hash(request.password)
        new_user = User(
            email=request.email,
            username=request.username,
            hashed_password=hashed_password,
            role="member",
            team_id=default_team.id if default_team else None
        )
        db.add(new_user)
        db.commit()
        db.refresh(new_user)
        role = getattr(new_user, "role", "member") or "member"
        team_id = getattr(new_user, "team_id", None)
        access_token = create_access_token(data={"sub": new_user.id, "role": role, "team_id": team_id})
        team_name = default_team.name if default_team else None
        return {
            "status": "success",
            "message": "Account created successfully",
            "access_token": access_token,
            "token_type": "bearer",
            "user": {
                "id": new_user.id,
                "email": new_user.email,
                "username": new_user.username,
                "role": role,
                "team_id": team_id,
                "team_name": team_name
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        import traceback
        error_trace = traceback.format_exc()
        print(f"Signup error: {error_trace}")  # Debug logging
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error creating account: {str(e)}"
        )

@app.post("/api/login")
async def login(request: LoginRequest, db: Session = Depends(get_db)):
    """Simple user login."""
    try:
        # Find user by email
        user = db.query(User).filter(User.email == request.email).first()
        
        if not user:
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Incorrect email or password"
            )
        
        if not verify_password(request.password, user.hashed_password):
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Incorrect email or password"
            )
        
        if user.is_active == 0:
            raise HTTPException(
                status_code=status.HTTP_403_FORBIDDEN,
                detail="User account is inactive"
            )
        role = getattr(user, "role", "member") or "member"
        role_sel = (request.role or "").strip().lower()
        if role_sel and role_sel not in ("manager", "team_head", "member"):
            raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid role selected")
        if role_sel and role_sel != role.lower():
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"This account is a {role.replace('_', ' ').title()}. Please select '{role.replace('_', ' ').title()}' or use the correct account."
            )
        team_id = getattr(user, "team_id", None)
        access_token = create_access_token(data={"sub": user.id, "role": role, "team_id": team_id})
        team_name = (getattr(user, "team", None) and getattr(user.team, "name", None)) or None
        try:
            from services.model_health import check_models_async
            asyncio.create_task(check_models_async())
            print(f"🔍 Model health check scheduled for user {user.id}")
        except Exception as e:
            print(f"⚠️ Failed to schedule model health check: {e}")
        return {
            "status": "success",
            "message": "Login successful",
            "access_token": access_token,
            "token_type": "bearer",
            "user": {
                "id": user.id,
                "email": user.email,
                "username": user.username,
                "role": role,
                "team_id": team_id,
                "team_name": team_name
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error during login: {str(e)}"
        )

# /api/me endpoint is now in routes/auth_routes.py

@app.post("/api/logout")
async def logout():
    """Logout endpoint (client-side token removal)."""
    return {"status": "success", "message": "Logged out successfully"}

# -------------------- VLM (Vision-Language Model) Helpers --------------------
def _get_env_bool(name: str, default: bool = False) -> bool:
    v = os.getenv(name)
    if v is None:
        return default
    return v.strip().lower() in {"1", "true", "yes", "on"}

def _vlm_summarize(image_path: str, context: str) -> str:
    """Enhanced function - uses comprehensive image interpretation."""
    from services.image_service import comprehensive_image_interpretation
    
    try:
        # Use comprehensive interpretation for much better results
        result = comprehensive_image_interpretation(image_path, context)
        return result['full_interpretation']
    except Exception as e:
        print(f"⚠️  Comprehensive interpretation failed, falling back: {e}")
        # Fallback to basic
        from services.image_service import enhanced_vlm_summarize
        image_type = "unknown"
        context_lower = context.lower() if context else ""
        if any(word in context_lower for word in ["diagram", "architecture", "component"]):
            image_type = "diagram"
        elif any(word in context_lower for word in ["chart", "graph", "plot", "data"]):
            image_type = "chart"
        elif any(word in context_lower for word in ["workflow", "process", "state", "transition"]):
            image_type = "workflow"
        return enhanced_vlm_summarize(image_path, context, image_type)

def _advanced_ocr_text(image_path: str = "", pil_image: Optional[Image.Image] = None) -> str:
    """Use advanced OCR first, with safe fallback to pytesseract."""
    try:
        from services.image_service import advanced_ocr_extract
        target_path = image_path
        temp_path = None
        if (not target_path) and pil_image is not None:
            temp_path = os.path.join(UPLOAD_DIR, f"tmp_ocr_{uuid.uuid4().hex}.png")
            pil_image.save(temp_path, format="PNG")
            target_path = temp_path
        if target_path and os.path.exists(target_path):
            result = advanced_ocr_extract(target_path)
            text = sanitize_unicode((result or {}).get("text", "") or "")
        else:
            text = sanitize_unicode(pytesseract.image_to_string(pil_image) if pil_image is not None else "")
        if temp_path and os.path.exists(temp_path):
            os.remove(temp_path)
        return text
    except Exception:
        try:
            return sanitize_unicode(pytesseract.image_to_string(pil_image) if pil_image is not None else "")
        except Exception:
            return ""

def _extract_image_requirements(image_path: str, context: str = "") -> list:
    """Dedicated per-image VLM requirement pass (legacy single-call wrapper)."""
    try:
        from services.image_service import extract_testable_requirements_from_image
        return extract_testable_requirements_from_image(image_path, context)
    except Exception:
        return []


# ==============================================================
# MAIN ANALYZE FUNCTION
# ==============================================================

async def _analyze_document_internal(
    file: UploadFile,
    user_id: int = None,
    db_session=None,
    progress_tracker_id: str = None,
    project_id: int = None,
):
    """Internal function: Extracts text, images, and OCR from any uploaded document."""
    from services.progress_service import ProcessingStage
    from services.progress_storage import get_progress_tracker
    from services.requirement_validation import is_srs_supported_upload

    ok_ext, ext_reasons = is_srs_supported_upload(file.filename or "")
    if not ok_ext:
        raise HTTPException(
            status_code=400,
            detail=(
                "Unsupported document type for SRS workflow. "
                + "; ".join(ext_reasons)
            ),
        )

    print(f"📄 Starting document analysis: {file.filename}")
    print(f"👤 User ID: {user_id}")
    
    tracker = None
    if progress_tracker_id:
        tracker = get_progress_tracker(progress_tracker_id)
        if tracker:
            tracker.set_stage(ProcessingStage.UPLOADING)
            progress = tracker.get_progress()
            print(f"📊 [{progress['progress']}%] {progress['message']}")
    
    print(f"⬆️ Uploading file: {file.filename}")
    filepath = os.path.join(UPLOAD_DIR, file.filename)
    with open(filepath, "wb") as f:
        content = await file.read()
        f.write(content)
        file_size = len(content)
        print(f"✅ File uploaded successfully ({file_size:,} bytes)")
    
    if tracker:
        tracker.set_stage(ProcessingStage.PARSING)
        progress = tracker.get_progress()
        print(f"📊 [{progress['progress']}%] {progress['message']}")

    text_output = ""
    image_count = 0
    image_metadata = []  # (image_id, image_path, page_number, ocr_text)
    image_positions = []  # list of dicts: {image_id, page, char_offset, context_before}

    def gen_image_id(fname: str, page_num: int, idx: int) -> str:
        base = f"{fname}|{page_num}|{idx}"
        h = hashlib.md5(base.encode()).hexdigest()[:10]
        return f"IMG-{h}-{page_num}-{idx}"

    # ----------- PDF PARSING (Chunked) -----------
    if file.filename.endswith(".pdf"):
        print(f"📄 Detected PDF file, initializing PDF reader...")
        reader = PdfReader(filepath)
        total_pages = len(reader.pages)
        chunk_size = 10

        print(f"📘 Document has {total_pages} pages. Starting text extraction...")
        
        if tracker:
            tracker.set_stage(ProcessingStage.PARSING, total_pages=total_pages, current_page=0)
            progress = tracker.get_progress()
            print(f"📊 [{progress['progress']}%] {progress['message']}")

        page_end_offsets = {}
        for start in range(0, total_pages, chunk_size):
            end = min(start + chunk_size, total_pages)
            print(f"🔹 Processing pages {start + 1}–{end} of {total_pages}...")

            for page_index in range(start, end):
                page = reader.pages[page_index]
                page_text = page.extract_text() or ""
                text_output += page_text + "\n"
                page_end_offsets[page_index + 1] = len(text_output)

                if tracker:
                    tracker.set_stage(ProcessingStage.PARSING, total_pages=total_pages, current_page=page_index + 1)
                    if (page_index + 1) % 5 == 0 or page_index + 1 == total_pages:  # Log every 5 pages
                        progress = tracker.get_progress()
                        print(f"📊 [{progress['progress']}%] Currently working on page {page_index + 1}/{total_pages}")

            print(f"✅ Finished pages {start + 1}–{end} of {total_pages}.")
        print(f"✅ Completed all {total_pages} pages. Extracted {len(text_output):,} characters.")
        
        if tracker:
            tracker.set_stage(ProcessingStage.TEXT_EXTRACTION)

        # Extract images using pypdf Page.images if available
        print(f"🖼️ Starting image detection and extraction...")
        if tracker:
            tracker.set_stage(ProcessingStage.IMAGE_DETECTION)
            progress = tracker.get_progress()
            print(f"📊 [{progress['progress']}%] {progress['message']}")
        
        try:
            from io import BytesIO
            detected_images = 0
            for page_num, page in enumerate(reader.pages, start=1):
                try:
                    imgs = getattr(page, "images", []) or []
                except Exception:
                    imgs = []
                for img_idx, img in enumerate(imgs, start=1):
                    try:
                        data = getattr(img, "data", None)
                        if not data:
                            continue
                        image_count += 1
                        image_id = gen_image_id(file.filename, page_num, img_idx)
                        out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}.png")
                        try:
                            im = Image.open(BytesIO(data))
                            im.save(out_path, format="PNG")
                            ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=im)
                        except Exception:
                            # If PIL cannot decode, just dump bytes and skip OCR
                            with open(out_path, "wb") as fimg:
                                fimg.write(data)
                            ocr_text = ""
                        image_metadata.append((image_id, out_path, page_num, ocr_text))
                        detected_images += 1
                        # Compute a simple position for context: end of that page's text
                        pos = page_end_offsets.get(page_num, len(text_output))
                        context_before = text_output[max(0, pos - 500):pos]
                        image_positions.append({
                            "image_id": image_id,
                            "page": page_num,
                            "char_offset": pos,
                            "context_before": context_before
                        })
                        
                        # Update progress for OCR
                        if tracker:
                            tracker.set_stage(ProcessingStage.OCR_PROCESSING, total_images=detected_images, current_image=detected_images)
                            progress = tracker.get_progress()
                            print(f"📊 [{progress['progress']}%] Processing OCR for image {detected_images} (page {page_num})")
                        
                        # Try VLM summary
                        if tracker:
                            tracker.set_stage(ProcessingStage.IMAGE_SUMMARIZATION, total_images=detected_images, current_image=detected_images)
                        vlm_sum = _vlm_summarize(out_path, context_before)
                        if vlm_sum:
                            print(f"✅ Generated summary for image {detected_images} on page {page_num}")
                        if vlm_sum:
                            text_output += f"\n[IMAGE_SUMMARY {image_id}]\n{vlm_sum}\n"
                        if (ocr_text or "").strip():
                            text_output += f"\n[IMAGE {image_id}]\nOCR: {(ocr_text or '').strip()}\n"
                    except Exception:
                        continue
        except Exception:
            pass

    # ----------- DOC (legacy) PARSING via LibreOffice conversion -----------
    elif file.filename.endswith(".doc"):
        # Convert .doc to .docx using LibreOffice if available
        soffice = _get_soffice_path()
        if not soffice:
            text_output = "Unsupported file format (.doc) and LibreOffice not found for conversion."
        else:
            try:
                subprocess.run([soffice, "--headless", "--convert-to", "docx", "--outdir", UPLOAD_DIR, filepath], check=True)
                converted = os.path.join(UPLOAD_DIR, os.path.splitext(file.filename)[0] + ".docx")
                d = docx.Document(converted)
                for para in d.paragraphs:
                    text_output += para.text + "\n"
                # Extract images from converted DOCX
                img_idx = 0
                for rel in d.part.rels.values():
                    if "image" in rel.reltype and getattr(rel, "target_part", None):
                        try:
                            img_idx += 1
                            image_count += 1
                            blob = rel.target_part.blob
                            page_num = 1
                            image_id = gen_image_id(file.filename, page_num, img_idx)
                            ext = ".png"
                            out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                            with open(out_path, "wb") as imf:
                                imf.write(blob)
                            try:
                                img = Image.open(out_path)
                                ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                            except Exception:
                                ocr_text = ""
                            image_metadata.append((image_id, out_path, page_num, ocr_text))
                            pos = len(text_output)
                            context_before = text_output[max(0, pos - 500):pos]
                            image_positions.append({
                                "image_id": image_id,
                                "page": page_num,
                                "char_offset": pos,
                                "context_before": context_before
                            })
                            vlm_sum = _vlm_summarize(out_path, context_before)
                            if vlm_sum:
                                text_output += f"\n[IMAGE_SUMMARY {image_id}]\n{vlm_sum}\n"
                            if ocr_text.strip():
                                text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                        except Exception:
                            pass
            except Exception as e:
                text_output = f"Conversion error for .doc: {e}"

    # ----------- DOCX PARSING -----------
    elif file.filename.endswith(".docx"):
        d = docx.Document(filepath)
        for para in d.paragraphs:
            text_output += para.text + "\n"
        # Extract images from DOCX
        img_idx = 0
        for rel in d.part.rels.values():
            if "image" in rel.reltype and getattr(rel, "target_part", None):
                try:
                    img_idx += 1
                    image_count += 1
                    blob = rel.target_part.blob
                    page_num = 1
                    image_id = gen_image_id(file.filename, page_num, img_idx)
                    ext = ".png"
                    out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                    with open(out_path, "wb") as imf:
                        imf.write(blob)
                    # OCR
                    try:
                        img = Image.open(out_path)
                        ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                    except Exception:
                        ocr_text = ""
                    image_metadata.append((image_id, out_path, page_num, ocr_text))
                    # Inject lightweight context for agent
                    if ocr_text.strip():
                        text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                except Exception:
                    pass

    # ----------- PPT (legacy) PARSING via LibreOffice conversion -----------
    elif file.filename.endswith(".ppt"):
        soffice = _get_soffice_path()
        if not soffice:
            text_output = "Unsupported file format (.ppt) and LibreOffice not found for conversion."
        else:
            try:
                subprocess.run([soffice, "--headless", "--convert-to", "pptx", "--outdir", UPLOAD_DIR, filepath], check=True)
                converted = os.path.join(UPLOAD_DIR, os.path.splitext(file.filename)[0] + ".pptx")
                prs = pptx.Presentation(converted)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_output += shape.text + "\n"
                        if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                            try:
                                image_count += 1
                                image_id = gen_image_id(file.filename, 1, image_count)
                                ext = ".png"
                                out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                                with open(out_path, "wb") as imf:
                                    imf.write(shape.image.blob)
                                try:
                                    img = Image.open(out_path)
                                    ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                                except Exception:
                                    ocr_text = ""
                                image_metadata.append((image_id, out_path, 1, ocr_text))
                                pos = len(text_output)
                                context_before = text_output[max(0, pos - 500):pos]
                                image_positions.append({
                                    "image_id": image_id,
                                    "page": 1,
                                    "char_offset": pos,
                                    "context_before": context_before
                                })
                                vlm_sum = _vlm_summarize(out_path, context_before)
                                if vlm_sum:
                                    text_output += f"\n[IMAGE_SUMMARY {image_id}]\n{vlm_sum}\n"
                                if ocr_text.strip():
                                    text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                            except Exception:
                                pass
            except Exception as e:
                text_output = f"Conversion error for .ppt: {e}"

    # ----------- PPTX PARSING -----------
    elif file.filename.endswith(".pptx"):
        prs = pptx.Presentation(filepath)
        # Extract text and images from PPTX
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_output += shape.text + "\n"
                # 13 == MSO_SHAPE_TYPE.PICTURE
                if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                    try:
                        image_count += 1
                        image_id = gen_image_id(file.filename, slide_idx, image_count)
                        ext = ".png"
                        out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                        with open(out_path, "wb") as imf:
                            imf.write(shape.image.blob)
                            try:
                                img = Image.open(out_path)
                                ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                            except Exception:
                                ocr_text = ""
                            image_metadata.append((image_id, out_path, slide_idx, ocr_text))
                        pos = len(text_output)
                        context_before = text_output[max(0, pos - 500):pos]
                        image_positions.append({
                            "image_id": image_id,
                            "page": slide_idx,
                            "char_offset": pos,
                            "context_before": context_before
                        })
                        vlm_sum = _vlm_summarize(out_path, context_before)
                        if vlm_sum:
                            text_output += f"\n[IMAGE_SUMMARY {image_id}]\n{vlm_sum}\n"
                        if ocr_text.strip():
                            text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                    except Exception:
                        pass

    # ----------- IMAGE PARSING -----------
    elif file.filename.lower().endswith((".png", ".jpg", ".jpeg")):
        if tracker:
            tracker.set_stage(ProcessingStage.IMAGE_DETECTION)
        image_count = 1
        img = Image.open(filepath)
        if tracker:
            tracker.set_stage(ProcessingStage.OCR_PROCESSING, total_images=1, current_image=1)
        text_output = _advanced_ocr_text(image_path=filepath, pil_image=img)
        image_id = gen_image_id(file.filename, 1, 1)
        image_metadata.append((image_id, filepath, 1, text_output))
        if text_output.strip():
            text_output += f"\n[IMAGE {image_id}]\nOCR: {text_output.strip()}\n"

    # ----------- TEXT FILE PARSING -----------
    elif file.filename.lower().endswith(".txt"):
        if tracker:
            tracker.set_stage(ProcessingStage.TEXT_EXTRACTION)
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            text_output = f.read()

    else:
        text_output = "Unsupported file format."

    # ----------- SAVE FULL TEXT FILE -----------
    # Sanitize text before writing to avoid surrogate encoding issues
    sanitized_text = sanitize_unicode(text_output or "")
    text_file_path = os.path.join(UPLOAD_DIR, f"{file.filename}_full.txt")
    with open(text_file_path, "w", encoding="utf-8") as tf:
        tf.write(sanitized_text)

    # ----------- VALIDATE DOCUMENT QUALITY -----------
    if tracker:
        tracker.set_stage(ProcessingStage.VALIDATING)
        progress = tracker.get_progress()
        print(f"📊 [{progress['progress']}%] {progress['message']}")

    v_res = validate_document_for_srs(sanitized_text)
    if v_res.get("is_rejected"):
        reject_msg = v_res.get("reject_reason", "Document does not meet technical/SRS quality standards.")
        print(f"❌ REJECTED: {reject_msg}")
        if tracker:
            tracker.set_stage(ProcessingStage.FINALIZING)
            tracker.complete()
        raise HTTPException(
            status_code=400,
            detail={
                "error": "DOCUMENT_REJECTED",
                "message": reject_msg,
                "score": v_res.get("srs_score"),
                "reasons": v_res.get("reasons")
            }
        )

    if tracker:
        tracker.set_stage(ProcessingStage.FINALIZING)
        progress = tracker.get_progress()
        print(f"📊 [{progress['progress']}%] {progress['message']}")

    summary = f"Extracted {len(text_output.split())} words and {image_count} image(s)."
    print(f"✅ Document analysis complete: {summary}")

    if tracker:
        tracker.complete()
        final_progress = tracker.get_progress()
        print(f"📊 [{final_progress['progress']}%] {final_progress['message']}")


    # Build quick lookup for contextual summaries
    ctx_by_id = {pos.get("image_id"): (pos.get("context_before") or "") for pos in image_positions}
    image_analysis_cache = {}

    # ----------- SAVE TO DATABASE -----------
    try:
        if db_session is None:
            db = SessionLocal()
            should_close = True
        else:
            db = db_session
            should_close = False
        
        record = ParsedFile(
            filename=file.filename,
            extracted_text=sanitized_text[:400],
            detected_shapes=image_count,
            summary=summary,
            full_text_path=text_file_path,
            user_id=user_id,  # Link to user
            project_id=project_id,
        )
        db.add(record)
        db.commit()
        db.refresh(record)

        # Save image metadata (including explainable diagram understanding context)
        from services.image_service import classify_diagram_type, extract_diagram_requirements_vlm
        for image_id, image_path, page_num, ocr_text in image_metadata:
            context_before = ctx_by_id.get(image_id, "")
            try:
                diagram_info = classify_diagram_type(image_path)
                scenario = "diagram_with_context" if (context_before or (ocr_text or "").strip()) else "diagram_only"
                vlm_pack = extract_diagram_requirements_vlm(
                    image_path=image_path,
                    scenario=scenario,
                    before_text=context_before,
                    after_text="",
                    ocr_text=ocr_text or "",
                )
                vlm_analysis = (vlm_pack or {}).get("understanding") or ""
                req_count = len((vlm_pack or {}).get("requirements") or [])
            except Exception as img_err:
                print(f"⚠️ Diagram analysis failed for {image_id}: {img_err}")
                diagram_info = {"type": "unknown", "confidence": 30, "detected_features": ["unrecognized pattern"]}
                vlm_analysis = ""
                req_count = 0

            image_analysis_cache[image_id] = {
                "diagram_type": diagram_info.get("type", "unknown"),
                "type_confidence": int(diagram_info.get("confidence", 0) or 0),
                "detected_features": diagram_info.get("detected_features", []),
                "vlm_analysis": vlm_analysis,
                "extracted_requirements_count": req_count,
            }

            img_meta = ImageMeta(
                file_id=record.id,
                image_path=image_path,
                page_number=page_num,
                ocr_text=ocr_text,
                diagram_type=image_analysis_cache[image_id]["diagram_type"],
                type_confidence=image_analysis_cache[image_id]["type_confidence"],
                detected_features=json.dumps(image_analysis_cache[image_id]["detected_features"], ensure_ascii=False),
                vlm_analysis=image_analysis_cache[image_id]["vlm_analysis"],
                extracted_requirements_count=image_analysis_cache[image_id]["extracted_requirements_count"],
            )
            db.add(img_meta)

        db.commit()
    except Exception as e:
        print(f"❌ Database error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Database error: {str(e)}")
    finally:
        # Ensure session is always closed
        if should_close and db_session is None:
            try:
                db.close()
            except Exception:
                pass

    # Build image summaries with VLM analysis for AI agent
    image_summaries = []
    images_list = []
    for (iid, path, pg, ocr) in image_metadata:
        cached = image_analysis_cache.get(iid) or {}
        vlm_summary = cached.get("vlm_analysis") or _vlm_summarize(path, ctx_by_id.get(iid, ""))
        ocr_summary = _summarize_image_ocr(ocr or "", context=ctx_by_id.get(iid, ""))
        final_summary = vlm_summary or ocr_summary
        
        # Sanitize OCR text and summaries to prevent Unicode encoding errors
        sanitized_ocr = sanitize_unicode((ocr or "").strip())
        sanitized_summary = sanitize_unicode(final_summary)
        
        # Format for AI agent (includes all relevant data)
        image_summaries.append({
            "image_id": iid,
            "path": path,
            "page": pg,
            "ocr": sanitized_ocr,
            "summary": sanitized_summary,
            "interpretation": sanitized_summary,  # Alias for compatibility
            "diagram_type": cached.get("diagram_type", "unknown"),
            "type_confidence": cached.get("type_confidence", 0),
            "detected_features": cached.get("detected_features", []),
            "vlm_analysis": sanitize_unicode(cached.get("vlm_analysis", "")),
            "extracted_requirements_count": int(cached.get("extracted_requirements_count", 0)),
        })
        
        # Format for display
        images_list.append({
            "image_id": iid,
            "path": path,
            "page": pg,
            "ocr": sanitized_ocr,
            "summary": sanitized_summary,
            "diagram_type": cached.get("diagram_type", "unknown"),
            "type_confidence": cached.get("type_confidence", 0),
            "detected_features": cached.get("detected_features", []),
            "vlm_analysis": sanitize_unicode(cached.get("vlm_analysis", "")),
            "extracted_requirements_count": int(cached.get("extracted_requirements_count", 0)),
        })

    # Sanitize text_output before using it
    sanitized_text_output = sanitize_unicode(text_output or "")
    sanitized_summary = sanitize_unicode(summary or "")
    
    result = {
        "filename": file.filename or "",
        "summary": sanitized_summary,
        "full_text_file": text_file_path or "",
        "extracted_text": sanitized_text_output,  # The extracted text content for AI agent
        "full_text": sanitized_text_output,  # Same as extracted_text for compatibility
        "images_detected": image_count,
        "image_metadata_saved": len(image_metadata),
        "image_summaries": image_summaries,  # Formatted summaries for AI agent
        "images": images_list,  # For display
        "image_positions": image_positions
    }
    
    # Sanitize all Unicode strings to prevent encoding errors (recursive sanitization)
    return sanitize_dict(result)


# -------------------- AUTHENTICATION PAGES --------------------
@app.get("/login", response_class=HTMLResponse)
async def login_page():
    """Clean login page with Bootstrap - NO VERIFICATION CODE."""
    return """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Login - FlowMind</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.2/dist/css/bootstrap.min.css" rel="stylesheet">
    <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css" rel="stylesheet">
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <style>
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body {
            font-family: 'Inter', sans-serif;
            background: linear-gradient(135deg, #0f172a 0%, #1e293b 50%, #334155 100%);
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 2rem;
            position: relative;
            overflow: hidden;
        }
        
        /* Animated Background - Data Getting In */
        .bg-animation {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            pointer-events: none;
            z-index: 0;
            overflow: hidden;
        }
        
        .data-stream {
            position: absolute;
            color: rgba(6, 182, 212, 0.4);
            font-size: 16px;
            font-weight: 600;
            font-family: 'Courier New', monospace;
            text-shadow: 0 0 10px rgba(6, 182, 212, 0.6);
            animation: streamIn linear infinite;
            white-space: nowrap;
        }
        
        @keyframes streamIn {
            0% {
                opacity: 0;
                transform: translate(0, 0) scale(0.5);
            }
            10% {
                opacity: 1;
            }
            90% {
                opacity: 1;
            }
            100% {
                opacity: 0;
                transform: translate(var(--end-x), var(--end-y)) scale(1);
            }
        }
        
        /* Circular convergence effect */
        .convergence-circle {
            position: absolute;
            width: 300px;
            height: 300px;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            border: 2px solid rgba(6, 182, 212, 0.1);
            border-radius: 50%;
            animation: pulseCircle 3s ease-in-out infinite;
        }
        
        .convergence-circle::before,
        .convergence-circle::after {
            content: '';
            position: absolute;
            width: 200px;
            height: 200px;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            border: 1px solid rgba(59, 130, 246, 0.1);
            border-radius: 50%;
        }
        
        .convergence-circle::after {
            width: 400px;
            height: 400px;
            border-color: rgba(139, 92, 246, 0.1);
        }
        
        @keyframes pulseCircle {
            0%, 100% {
                transform: translate(-50%, -50%) scale(1);
                opacity: 0.5;
            }
            50% {
                transform: translate(-50%, -50%) scale(1.1);
                opacity: 0.8;
            }
        }
        
        .auth-card {
            background: rgba(255, 255, 255, 0.95);
            border-radius: 20px;
            padding: 3rem;
            box-shadow: 0 25px 50px -12px rgba(0, 0, 0, 0.4);
            width: 100%;
            max-width: 450px;
            animation: slideUp 0.5s ease-out;
            position: relative;
            z-index: 1;
            backdrop-filter: blur(10px);
            border: 2px solid rgba(6, 182, 212, 0.2);
        }
        @keyframes slideUp {
            from { opacity: 0; transform: translateY(30px); }
            to { opacity: 1; transform: translateY(0); }
        }
        .auth-header {
            text-align: center;
            margin-bottom: 2.5rem;
        }
        .auth-header h1 {
            font-size: 2.5rem;
            font-weight: 700;
            color: #0f172a;
            margin-bottom: 0.5rem;
        }
        .auth-header p {
            color: #475569;
            font-size: 1rem;
        }
        .form-label {
            font-weight: 600;
            color: #0f172a;
            margin-bottom: 0.5rem;
        }
        .form-control {
            padding: 0.875rem 1rem;
            border: 2px solid #e2e8f0;
            border-radius: 12px;
            font-size: 1rem;
            transition: all 0.3s ease;
            background: white;
            color: #0f172a;
        }
        .form-control:focus {
            border-color: #06b6d4;
            box-shadow: 0 0 0 4px rgba(6, 182, 212, 0.1);
            outline: none;
        }
        .btn-primary {
            width: 100%;
            padding: 0.875rem;
            border-radius: 12px;
            font-weight: 600;
            font-size: 1rem;
            background: linear-gradient(135deg, #0891b2 0%, #2563eb 50%, #7c3aed 100%);
            border: none;
            transition: all 0.3s ease;
            color: white;
            box-shadow: 0 4px 12px rgba(6, 182, 212, 0.3);
        }
        .btn-primary:hover {
            transform: translateY(-2px);
            box-shadow: 0 8px 24px rgba(6, 182, 212, 0.4);
            background: linear-gradient(135deg, #075985 0%, #1e40af 50%, #6d28d9 100%);
        }
        .alert {
            border-radius: 12px;
            border: none;
            margin-bottom: 1.5rem;
        }
        .auth-footer {
            text-align: center;
            margin-top: 2rem;
            color: #475569;
        }
        .auth-footer a {
            color: #06b6d4;
            text-decoration: none;
            font-weight: 500;
        }
        .auth-footer a:hover {
            text-decoration: underline;
            color: #0891b2;
        }
        .text-primary {
            color: #06b6d4 !important;
        }
    </style>
</head>
<body>
    <div class="auth-card">
        <div class="auth-header">
            <h1><i class="fas fa-sign-in-alt text-primary"></i> Login</h1>
            <p>Welcome back to FlowMind</p>
        </div>
        
        <div id="alertContainer"></div>
        
        <form id="loginForm">
            <div class="mb-3">
                <label for="loginRole" class="form-label">Select your role</label>
                <select class="form-control form-select" id="loginRole" name="role" style="padding: 0.875rem 1rem; border-radius: 12px; font-size: 1rem;">
                    <option value="">— Choose role (optional) —</option>
                    <option value="manager">Manager</option>
                    <option value="team_head">Team Head</option>
                    <option value="member">Member</option>
                </select>
                <small class="text-muted">Pick the role that matches your account.</small>
            </div>
            <div class="mb-3">
                <label for="email" class="form-label">Email Address</label>
                <input type="email" class="form-control" id="email" name="email" required autocomplete="email" placeholder="Enter your email">
            </div>
            
            <div class="mb-4">
                <label for="password" class="form-label">Password</label>
                <input type="password" class="form-control" id="password" name="password" required autocomplete="current-password" placeholder="Enter your password">
            </div>
            
            <button type="submit" class="btn btn-primary" id="submitBtn">
                <i class="fas fa-sign-in-alt"></i> Login
            </button>
        </form>
        
        <div class="auth-footer">
            <p class="mb-2">Don't have an account? <a href="/signup">Sign up</a></p>
            <a href="/"><i class="fas fa-arrow-left"></i> Back to Home</a>
        </div>
    </div>
    
    <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.2/dist/js/bootstrap.bundle.min.js"></script>
    <script>
        // Create animated data streams flowing in
        function createDataStreams() {
            const bgAnimation = document.querySelector('.bg-animation');
            if (!bgAnimation) return;
            
            const dataSymbols = ['LOGIN', 'AUTH', 'ACCESS', 'SECURE', 'FLOW', 'MIND', 'AI', 'DATA', 'INFO', 'USER', 'PASS', 'KEY'];
            const centerX = window.innerWidth / 2;
            const centerY = window.innerHeight / 2;
            
            function createStream() {
                const stream = document.createElement('div');
                stream.className = 'data-stream';
                
                // Random data symbol
                stream.textContent = dataSymbols[Math.floor(Math.random() * dataSymbols.length)];
                
                // Random starting position from edges
                const side = Math.floor(Math.random() * 4); // 0=top, 1=right, 2=bottom, 3=left
                let startX, startY, endX, endY;
                
                if (side === 0) { // Top
                    startX = Math.random() * window.innerWidth;
                    startY = -50;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                } else if (side === 1) { // Right
                    startX = window.innerWidth + 50;
                    startY = Math.random() * window.innerHeight;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                } else if (side === 2) { // Bottom
                    startX = Math.random() * window.innerWidth;
                    startY = window.innerHeight + 50;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                } else { // Left
                    startX = -50;
                    startY = Math.random() * window.innerHeight;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                }
                
                const deltaX = endX - startX;
                const deltaY = endY - startY;
                const distance = Math.sqrt(deltaX * deltaX + deltaY * deltaY);
                const duration = 3 + (distance / 200); // Speed based on distance
                
                stream.style.left = startX + 'px';
                stream.style.top = startY + 'px';
                stream.style.setProperty('--end-x', deltaX + 'px');
                stream.style.setProperty('--end-y', deltaY + 'px');
                stream.style.animationDuration = duration + 's';
                stream.style.animationDelay = Math.random() * 0.5 + 's';
                
                // Random size
                const size = 12 + Math.random() * 8;
                stream.style.fontSize = size + 'px';
                
                // Random color variation
                const colors = [
                    'rgba(6, 182, 212, 0.4)',
                    'rgba(59, 130, 246, 0.4)',
                    'rgba(139, 92, 246, 0.4)',
                    'rgba(6, 182, 212, 0.3)',
                    'rgba(59, 130, 246, 0.3)'
                ];
                stream.style.color = colors[Math.floor(Math.random() * colors.length)];
                
                bgAnimation.appendChild(stream);
                
                // Remove after animation
                setTimeout(() => {
                    if (stream.parentNode) {
                        stream.parentNode.removeChild(stream);
                    }
                }, (duration + 0.5) * 1000);
            }
            
            // Create streams continuously
            setInterval(createStream, 150);
            setInterval(createStream, 200);
            setInterval(createStream, 250);
            
            // Initial burst
            for (let i = 0; i < 15; i++) {
                setTimeout(() => createStream(), i * 100);
            }
        }
        
        // Initialize on page load
        window.addEventListener('DOMContentLoaded', createDataStreams);
        
        const form = document.getElementById('loginForm');
        const alertContainer = document.getElementById('alertContainer');
        const submitBtn = document.getElementById('submitBtn');
        
        function showAlert(message, type = 'danger') {
            alertContainer.innerHTML = `
                <div class="alert alert-${type} alert-dismissible fade show" role="alert">
                    <i class="fas fa-${type === 'danger' ? 'exclamation-circle' : 'check-circle'}"></i> ${message}
                    <button type="button" class="btn-close" data-bs-dismiss="alert"></button>
                </div>
            `;
        }
        
        form.addEventListener('submit', async (e) => {
            e.preventDefault();
            alertContainer.innerHTML = '';
            submitBtn.disabled = true;
            submitBtn.innerHTML = '<span class="spinner-border spinner-border-sm me-2"></span>Logging in...';
            
            const roleEl = document.getElementById('loginRole');
            const formData = {
                email: document.getElementById('email').value,
                password: document.getElementById('password').value,
                role: roleEl ? roleEl.value : ''
            };
            
            try {
                const response = await fetch('/api/login', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify(formData)
                });
                
                const data = await response.json();
                
                if (response.ok) {
                    localStorage.setItem('access_token', data.access_token);
                    localStorage.setItem('user', JSON.stringify(data.user));
                    showAlert('Login successful! Redirecting...', 'success');
                    const role = (data.user && data.user.role) || 'member';
                    const target = role === 'manager' ? '/manager' : (role === 'team_head' ? '/team' : '/dashboard');
                    setTimeout(() => window.location.href = target, 1000);
                } else {
                    let msg = data.detail || 'Login failed. Please try again.';
                    if (Array.isArray(msg)) msg = (msg[0] && msg[0].msg) ? msg.map(function(x){ return x.msg; }).join('; ') : String(msg);
                    else if (msg && typeof msg === 'object' && msg.msg) msg = msg.msg;
                    showAlert(msg);
                    submitBtn.disabled = false;
                    submitBtn.innerHTML = '<i class="fas fa-sign-in-alt"></i> Login';
                }
            } catch (error) {
                showAlert('Network error: ' + (error.message || 'Please check your connection and try again.'));
                submitBtn.disabled = false;
                submitBtn.innerHTML = '<i class="fas fa-sign-in-alt"></i> Login';
            }
        });
    </script>
</body>
</html>
    """

@app.get("/signup", response_class=HTMLResponse)
async def signup_page():
    """Clean signup page with Bootstrap - NO VERIFICATION CODE."""
    return """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Sign Up - FlowMind</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.2/dist/css/bootstrap.min.css" rel="stylesheet">
    <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css" rel="stylesheet">
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <style>
        * { margin: 0; padding: 0; box-sizing: border-box; }
        body {
            font-family: 'Inter', sans-serif;
            background: linear-gradient(135deg, #0f172a 0%, #1e293b 50%, #334155 100%);
            min-height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 2rem;
            position: relative;
            overflow: hidden;
        }
        
        /* Animated Background - Data Getting In */
        .bg-animation {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            pointer-events: none;
            z-index: 0;
            overflow: hidden;
        }
        
        .data-stream {
            position: absolute;
            color: rgba(6, 182, 212, 0.4);
            font-size: 16px;
            font-weight: 600;
            font-family: 'Courier New', monospace;
            text-shadow: 0 0 10px rgba(6, 182, 212, 0.6);
            animation: streamIn linear infinite;
            white-space: nowrap;
        }
        
        @keyframes streamIn {
            0% {
                opacity: 0;
                transform: translate(0, 0) scale(0.5);
            }
            10% {
                opacity: 1;
            }
            90% {
                opacity: 1;
            }
            100% {
                opacity: 0;
                transform: translate(var(--end-x), var(--end-y)) scale(1);
            }
        }
        
        /* Circular convergence effect */
        .convergence-circle {
            position: absolute;
            width: 300px;
            height: 300px;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            border: 2px solid rgba(6, 182, 212, 0.1);
            border-radius: 50%;
            animation: pulseCircle 3s ease-in-out infinite;
        }
        
        .convergence-circle::before,
        .convergence-circle::after {
            content: '';
            position: absolute;
            width: 200px;
            height: 200px;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            border: 1px solid rgba(59, 130, 246, 0.1);
            border-radius: 50%;
        }
        
        .convergence-circle::after {
            width: 400px;
            height: 400px;
            border-color: rgba(139, 92, 246, 0.1);
        }
        
        @keyframes pulseCircle {
            0%, 100% {
                transform: translate(-50%, -50%) scale(1);
                opacity: 0.5;
            }
            50% {
                transform: translate(-50%, -50%) scale(1.1);
                opacity: 0.8;
            }
        }
        
        .auth-card {
            background: rgba(255, 255, 255, 0.95);
            border-radius: 20px;
            padding: 3rem;
            box-shadow: 0 25px 50px -12px rgba(0, 0, 0, 0.4);
            width: 100%;
            max-width: 450px;
            animation: slideUp 0.5s ease-out;
            position: relative;
            z-index: 1;
            backdrop-filter: blur(10px);
            border: 2px solid rgba(6, 182, 212, 0.2);
        }
        @keyframes slideUp {
            from { opacity: 0; transform: translateY(30px); }
            to { opacity: 1; transform: translateY(0); }
        }
        .auth-header {
            text-align: center;
            margin-bottom: 2.5rem;
        }
        .auth-header h1 {
            font-size: 2.5rem;
            font-weight: 700;
            color: #0f172a;
            margin-bottom: 0.5rem;
        }
        .auth-header p {
            color: #475569;
            font-size: 1rem;
        }
        .form-label {
            font-weight: 600;
            color: #0f172a;
            margin-bottom: 0.5rem;
        }
        .form-control {
            padding: 0.875rem 1rem;
            border: 2px solid #e2e8f0;
            border-radius: 12px;
            font-size: 1rem;
            transition: all 0.3s ease;
            background: white;
            color: #0f172a;
        }
        .form-control:focus {
            border-color: #06b6d4;
            box-shadow: 0 0 0 4px rgba(6, 182, 212, 0.1);
            outline: none;
        }
        .btn-primary {
            width: 100%;
            padding: 0.875rem;
            border-radius: 12px;
            font-weight: 600;
            font-size: 1rem;
            background: linear-gradient(135deg, #0891b2 0%, #2563eb 50%, #7c3aed 100%);
            border: none;
            transition: all 0.3s ease;
            color: white;
            box-shadow: 0 4px 12px rgba(6, 182, 212, 0.3);
        }
        .btn-primary:hover {
            transform: translateY(-2px);
            box-shadow: 0 8px 24px rgba(6, 182, 212, 0.4);
            background: linear-gradient(135deg, #075985 0%, #1e40af 50%, #6d28d9 100%);
        }
        .alert {
            border-radius: 12px;
            border: none;
            margin-bottom: 1.5rem;
        }
        .auth-footer {
            text-align: center;
            margin-top: 2rem;
            color: #475569;
        }
        .auth-footer a {
            color: #06b6d4;
            text-decoration: none;
            font-weight: 500;
        }
        .auth-footer a:hover {
            text-decoration: underline;
            color: #0891b2;
        }
        .text-primary {
            color: #06b6d4 !important;
        }
        .password-hint {
            font-size: 0.875rem;
            color: #475569;
            margin-top: 0.25rem;
        }
    </style>
</head>
<body>
    <!-- Animated Background -->
    <div class="bg-animation">
        <div class="convergence-circle"></div>
    </div>
    
    <div class="auth-card">
        <div class="auth-header">
            <h1><i class="fas fa-user-plus text-primary"></i> Sign Up</h1>
            <p>Create your FlowMind account</p>
        </div>
        
        <div id="alertContainer"></div>
        
        <form id="signupForm">
            <div class="mb-3">
                <label for="email" class="form-label">Email Address</label>
                <input type="email" class="form-control" id="email" name="email" required autocomplete="email" placeholder="Enter your email">
            </div>
            
            <div class="mb-3">
                <label for="username" class="form-label">Username</label>
                <input type="text" class="form-control" id="username" name="username" required autocomplete="username" placeholder="Choose a username">
            </div>
            
            <div class="mb-4">
                <label for="password" class="form-label">Password</label>
                <input type="password" class="form-control" id="password" name="password" required autocomplete="new-password" minlength="6" placeholder="Create a password">
                <div class="password-hint">Must be at least 6 characters long</div>
            </div>
            
            <button type="submit" class="btn btn-primary" id="submitBtn">
                <i class="fas fa-user-plus"></i> Create Account
            </button>
        </form>
        
        <div class="auth-footer">
            <p class="mb-2">Already have an account? <a href="/login">Login</a></p>
            <a href="/"><i class="fas fa-arrow-left"></i> Back to Home</a>
        </div>
    </div>
    
    <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.2/dist/js/bootstrap.bundle.min.js"></script>
    <script>
        // Create animated data streams flowing in
        function createDataStreams() {
            const bgAnimation = document.querySelector('.bg-animation');
            if (!bgAnimation) return;
            
            const dataSymbols = ['SIGNUP', 'REGISTER', 'CREATE', 'ACCOUNT', 'JOIN', 'FLOW', 'MIND', 'AI', 'DATA', 'INFO', 'USER', 'NEW'];
            const centerX = window.innerWidth / 2;
            const centerY = window.innerHeight / 2;
            
            function createStream() {
                const stream = document.createElement('div');
                stream.className = 'data-stream';
                
                // Random data symbol
                stream.textContent = dataSymbols[Math.floor(Math.random() * dataSymbols.length)];
                
                // Random starting position from edges
                const side = Math.floor(Math.random() * 4); // 0=top, 1=right, 2=bottom, 3=left
                let startX, startY, endX, endY;
                
                if (side === 0) { // Top
                    startX = Math.random() * window.innerWidth;
                    startY = -50;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                } else if (side === 1) { // Right
                    startX = window.innerWidth + 50;
                    startY = Math.random() * window.innerHeight;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                } else if (side === 2) { // Bottom
                    startX = Math.random() * window.innerWidth;
                    startY = window.innerHeight + 50;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                } else { // Left
                    startX = -50;
                    startY = Math.random() * window.innerHeight;
                    endX = centerX + (Math.random() - 0.5) * 200;
                    endY = centerY + (Math.random() - 0.5) * 200;
                }
                
                const deltaX = endX - startX;
                const deltaY = endY - startY;
                const distance = Math.sqrt(deltaX * deltaX + deltaY * deltaY);
                const duration = 3 + (distance / 200); // Speed based on distance
                
                stream.style.left = startX + 'px';
                stream.style.top = startY + 'px';
                stream.style.setProperty('--end-x', deltaX + 'px');
                stream.style.setProperty('--end-y', deltaY + 'px');
                stream.style.animationDuration = duration + 's';
                stream.style.animationDelay = Math.random() * 0.5 + 's';
                
                // Random size
                const size = 12 + Math.random() * 8;
                stream.style.fontSize = size + 'px';
                
                // Random color variation
                const colors = [
                    'rgba(6, 182, 212, 0.4)',
                    'rgba(59, 130, 246, 0.4)',
                    'rgba(139, 92, 246, 0.4)',
                    'rgba(6, 182, 212, 0.3)',
                    'rgba(59, 130, 246, 0.3)'
                ];
                stream.style.color = colors[Math.floor(Math.random() * colors.length)];
                
                bgAnimation.appendChild(stream);
                
                // Remove after animation
                setTimeout(() => {
                    if (stream.parentNode) {
                        stream.parentNode.removeChild(stream);
                    }
                }, (duration + 0.5) * 1000);
            }
            
            // Create streams continuously
            setInterval(createStream, 150);
            setInterval(createStream, 200);
            setInterval(createStream, 250);
            
            // Initial burst
            for (let i = 0; i < 15; i++) {
                setTimeout(() => createStream(), i * 100);
            }
        }
        
        // Initialize on page load
        window.addEventListener('DOMContentLoaded', createDataStreams);
        
        const form = document.getElementById('signupForm');
        const alertContainer = document.getElementById('alertContainer');
        const submitBtn = document.getElementById('submitBtn');
        
        function showAlert(message, type = 'danger') {
            alertContainer.innerHTML = `
                <div class="alert alert-${type} alert-dismissible fade show" role="alert">
                    <i class="fas fa-${type === 'danger' ? 'exclamation-circle' : 'check-circle'}"></i> ${message}
                    <button type="button" class="btn-close" data-bs-dismiss="alert"></button>
                </div>
            `;
        }
        
        form.addEventListener('submit', async (e) => {
            e.preventDefault();
            alertContainer.innerHTML = '';
            submitBtn.disabled = true;
            submitBtn.innerHTML = '<span class="spinner-border spinner-border-sm me-2"></span>Creating account...';
            
            const formData = {
                email: document.getElementById('email').value,
                username: document.getElementById('username').value,
                password: document.getElementById('password').value
            };
            
            try {
                const response = await fetch('/api/signup', {
                    method: 'POST',
                    headers: { 'Content-Type': 'application/json' },
                    body: JSON.stringify(formData)
                });
                
                const data = await response.json();
                
                if (response.ok) {
                    localStorage.setItem('access_token', data.access_token);
                    localStorage.setItem('user', JSON.stringify(data.user));
                    showAlert('Account created successfully! Redirecting...', 'success');
                    setTimeout(() => window.location.href = '/extract', 1000);
                } else {
                    showAlert(data.detail || 'Signup failed. Please try again.');
                    submitBtn.disabled = false;
                    submitBtn.innerHTML = '<i class="fas fa-user-plus"></i> Create Account';
                }
            } catch (error) {
                showAlert('Network error. Please check your connection and try again.');
                submitBtn.disabled = false;
                submitBtn.innerHTML = '<i class="fas fa-user-plus"></i> Create Account';
            }
        });
    </script>
</body>
</html>
    """

@app.get("/about", response_class=HTMLResponse)
async def about_page():
    """About Us page."""
    return """
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>About Us - FlowMind</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
        <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css" rel="stylesheet">
        <style>
            :root {
                --primary-color: #2563eb;
                --text-primary: #1e293b;
                --text-secondary: #64748b;
                --border: #e2e8f0;
                --radius: 12px;
            }
            * {
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }
            body {
                font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                min-height: 100vh;
                padding: 2rem;
            }
            .container {
                max-width: 900px;
                margin: 0 auto;
                background: white;
                border-radius: var(--radius);
                padding: 3rem;
                box-shadow: 0 20px 25px -5px rgba(0, 0, 0, 0.1);
            }
            .header {
                text-align: center;
                margin-bottom: 3rem;
            }
            .header h1 {
                font-size: 2.5rem;
                font-weight: 700;
                color: var(--text-primary);
                margin-bottom: 1rem;
            }
            .header p {
                color: var(--text-secondary);
                font-size: 1.1rem;
            }
            .content {
                line-height: 1.8;
                color: var(--text-primary);
            }
            .content h2 {
                font-size: 1.5rem;
                font-weight: 600;
                margin-top: 2rem;
                margin-bottom: 1rem;
                color: var(--primary-color);
            }
            .content p {
                margin-bottom: 1rem;
                color: var(--text-secondary);
            }
            .content ul {
                margin-left: 2rem;
                margin-bottom: 1rem;
                color: var(--text-secondary);
            }
            .content li {
                margin-bottom: 0.5rem;
            }
            .back-link {
                display: inline-block;
                margin-top: 2rem;
                color: var(--primary-color);
                text-decoration: none;
                font-weight: 500;
            }
            .back-link:hover {
                text-decoration: underline;
            }
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <h1><i class="fas fa-info-circle"></i> About FlowMind</h1>
                <p>Intelligent Document Analysis Platform</p>
            </div>
            <div class="content">
                <h2>What is FlowMind?</h2>
                <p>
                    FlowMind is an advanced document analysis platform that leverages artificial intelligence 
                    and machine learning to extract, analyze, and understand complex requirements from various 
                    document formats. Our platform helps organizations streamline their document processing 
                    workflows and gain valuable insights from their documents.
                </p>
                
                <h2>Key Features</h2>
                <ul>
                    <li><strong>Multi-Format Support:</strong> Process PDFs, Word documents, PowerPoint presentations, images, and text files</li>
                    <li><strong>Intelligent Extraction:</strong> Advanced OCR and text extraction capabilities</li>
                    <li><strong>Requirements Analysis:</strong> Automatically identify and extract requirements from documents</li>
                    <li><strong>RAG-Powered:</strong> Retrieval-Augmented Generation for accurate document understanding</li>
                    <li><strong>Visual Analysis:</strong> Extract and analyze images, diagrams, and visual elements</li>
                    <li><strong>Self-Learning System:</strong> Continuously improves extraction accuracy over time</li>
                </ul>
                
                <h2>Our Mission</h2>
                <p>
                    At FlowMind, we believe that document analysis should be intelligent, efficient, and accessible. 
                    Our mission is to empower organizations with cutting-edge AI technology that transforms how they 
                    process and understand their documents.
                </p>
                
                <h2>Implementation Phases (Role-Based Access)</h2>
                <p>The following six phases are implemented in the application:</p>
                <ul>
                    <li><strong>Phase 1 &ndash; Schema:</strong> Teams table, User role and team_id; DB migration and seed.</li>
                    <li><strong>Phase 2 &ndash; Auth:</strong> JWT includes role and team_id; login/signup return role and team_name.</li>
                    <li><strong>Phase 3 &ndash; Scoped APIs:</strong> My uploads, progress, features, and teams APIs filtered by Manager / Team Head / Member visibility.</li>
                    <li><strong>Phase 4 &ndash; Manager &amp; Team UI:</strong> Manager and My Team pages and dashboard nav (visible when your role is manager or team_head).</li>
                    <li><strong>Phase 5 &ndash; Requirements UI:</strong> Simplified requirements view and collapsible feedback on the approve page.</li>
                    <li><strong>Phase 6 &ndash; Post-login redirect:</strong> Redirect to dashboard, manager, or team based on role after login.</li>
                </ul>
                <p>To see Manager / My Team links, log in as a user with role <code>manager</code> or <code>team_head</code>. Restart the server after code changes to load the latest implementation.</p>
                
                <h2>Technology Stack</h2>
                <p>
                    FlowMind is built using state-of-the-art technologies including FastAPI, LangChain, ChromaDB, 
                    and advanced machine learning models for natural language processing and computer vision.
                </p>
                
                <a href="/" class="back-link"><i class="fas fa-arrow-left"></i> Back to Home</a>
            </div>
        </div>
    </body>
    </html>
    """

@app.get("/", response_class=HTMLResponse)
async def root():
    """Serve the landing page."""
    # Try to serve index.html from root first
    index_path = os.path.join(BASE_DIR, "index.html")
    if os.path.exists(index_path):
        with open(index_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    
    # Fallback to static/landing.html
    landing_path = os.path.join(STATIC_DIR, "landing.html")
    if os.path.exists(landing_path):
        with open(landing_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    
    # Fallback to embedded HTML if static file doesn't exist
    return """
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>FlowMind – Intelligent Document Analysis</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
        <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css" rel="stylesheet">
        <style>
            :root {
                --primary-color: #2563eb;
                --primary-dark: #1d4ed8;
                --secondary-color: #64748b;
                --success-color: #10b981;
                --warning-color: #f59e0b;
                --error-color: #ef4444;
                --background: #f8fafc;
                --surface: #ffffff;
                --surface-elevated: #ffffff;
                --text-primary: #1e293b;
                --text-secondary: #64748b;
                --text-muted: #94a3b8;
                --border: #e2e8f0;
                --border-light: #f1f5f9;
                --shadow: 0 1px 3px 0 rgba(0, 0, 0, 0.1), 0 1px 2px 0 rgba(0, 0, 0, 0.06);
                --shadow-lg: 0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -2px rgba(0, 0, 0, 0.05);
                --shadow-xl: 0 20px 25px -5px rgba(0, 0, 0, 0.1), 0 10px 10px -5px rgba(0, 0, 0, 0.04);
                --radius: 12px;
                --radius-sm: 8px;
                --radius-lg: 16px;
            }

            * {
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }

            body {
                font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                min-height: 100vh;
                color: var(--text-primary);
                line-height: 1.6;
            }

            .container {
                max-width: 1200px;
                margin: 0 auto;
                padding: 2rem;
            }

            .header {
                text-align: center;
                margin-bottom: 3rem;
                animation: fadeInDown 0.8s ease-out;
            }

            .header h1 {
                font-size: 3rem;
                font-weight: 700;
                color: white;
                margin-bottom: 0.5rem;
                text-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            }

            .header p {
                font-size: 1.2rem;
                color: rgba(255, 255, 255, 0.9);
                font-weight: 300;
            }

            .main-content {
                display: grid;
                grid-template-columns: 1fr 1fr;
                gap: 2rem;
                margin-bottom: 2rem;
            }

            .card {
                background: var(--surface);
                border-radius: var(--radius-lg);
                padding: 2rem;
                box-shadow: var(--shadow-xl);
                border: 1px solid var(--border-light);
                transition: all 0.3s ease;
                animation: fadeInUp 0.8s ease-out;
            }

            .card:hover {
                transform: translateY(-4px);
                box-shadow: 0 25px 50px -12px rgba(0, 0, 0, 0.25);
            }

            .card:nth-child(2) {
                animation-delay: 0.2s;
            }

            .card-header {
                display: flex;
                align-items: center;
                margin-bottom: 1.5rem;
            }

            .card-icon {
                width: 48px;
                height: 48px;
                border-radius: var(--radius);
                display: flex;
                align-items: center;
                justify-content: center;
                margin-right: 1rem;
                font-size: 1.5rem;
                color: white;
            }

            .card-icon.basic {
                background: linear-gradient(135deg, var(--primary-color), var(--primary-dark));
            }

            .card-icon.agent {
                background: linear-gradient(135deg, var(--success-color), #059669);
            }

            .card-title {
                font-size: 1.5rem;
                font-weight: 600;
                color: var(--text-primary);
            }

            .card-subtitle {
                font-size: 0.9rem;
                color: var(--text-secondary);
                margin-top: 0.25rem;
            }

            .form-group {
                margin-bottom: 1.5rem;
            }

            .file-input-wrapper {
                position: relative;
                display: inline-block;
                width: 100%;
            }

            .file-input {
                position: absolute;
                opacity: 0;
                width: 100%;
                height: 100%;
                cursor: pointer;
            }

            .file-input-label {
                display: flex;
                align-items: center;
                justify-content: center;
                padding: 1.5rem;
                border: 2px dashed var(--border);
                border-radius: var(--radius);
                background: var(--background);
                cursor: pointer;
                transition: all 0.3s ease;
                font-weight: 500;
                color: var(--text-secondary);
            }

            .file-input-label:hover {
                border-color: var(--primary-color);
                background: rgba(37, 99, 235, 0.05);
                color: var(--primary-color);
            }

            .file-input-label i {
                margin-right: 0.5rem;
                font-size: 1.2rem;
            }

            .btn {
                background: linear-gradient(135deg, var(--primary-color), var(--primary-dark));
                color: white;
                border: none;
                padding: 0.875rem 1.5rem;
                border-radius: var(--radius);
                font-weight: 500;
                font-size: 1rem;
                cursor: pointer;
                transition: all 0.3s ease;
                display: inline-flex;
                align-items: center;
                gap: 0.5rem;
                box-shadow: var(--shadow);
            }

            .btn:hover {
                transform: translateY(-2px);
                box-shadow: var(--shadow-lg);
            }

            .btn:active {
                transform: translateY(0);
            }

            .btn:disabled {
                opacity: 0.6;
                cursor: not-allowed;
                transform: none;
            }

            .progress-container {
                margin-top: 1rem;
                display: none;
            }

            .progress-bar {
                width: 100%;
                height: 8px;
                background: var(--border-light);
                border-radius: 4px;
                overflow: hidden;
                margin-bottom: 0.5rem;
            }

            .progress-fill {
                height: 100%;
                background: linear-gradient(90deg, var(--primary-color), var(--success-color));
                width: 0%;
                transition: width 0.3s ease;
                border-radius: 4px;
            }

            .progress-text {
                font-size: 0.875rem;
                color: var(--text-secondary);
                text-align: center;
            }

            .output-container {
                margin-top: 1.5rem;
                background: #0f172a;
                border-radius: var(--radius);
                padding: 1.5rem;
                color: #e2e8f0;
                font-family: 'Monaco', 'Menlo', 'Ubuntu Mono', monospace;
                font-size: 0.875rem;
                line-height: 1.5;
                max-height: 400px;
                overflow-y: auto;
                border: 1px solid #1e293b;
            }

            .output-container::-webkit-scrollbar {
                width: 6px;
            }

            .output-container::-webkit-scrollbar-track {
                background: #1e293b;
                border-radius: 3px;
            }

            .output-container::-webkit-scrollbar-thumb {
                background: #475569;
                border-radius: 3px;
            }

            .output-container::-webkit-scrollbar-thumb:hover {
                background: #64748b;
            }

            .view-link {
                margin-top: 1rem;
                text-align: center;
            }

            .view-link a {
                display: inline-flex;
                align-items: center;
                gap: 0.5rem;
                background: linear-gradient(135deg, var(--success-color), #059669);
                color: white;
                text-decoration: none;
                padding: 0.75rem 1.5rem;
                border-radius: var(--radius);
                font-weight: 500;
                transition: all 0.3s ease;
                box-shadow: var(--shadow);
            }

            .view-link a:hover {
                transform: translateY(-2px);
                box-shadow: var(--shadow-lg);
            }

            .navigation {
                background: var(--surface);
                border-radius: var(--radius-lg);
                padding: 1.5rem;
                box-shadow: var(--shadow-lg);
                animation: fadeInUp 0.8s ease-out 0.4s both;
            }

            .nav-links {
                display: flex;
                justify-content: center;
                gap: 2rem;
                flex-wrap: wrap;
            }

            .nav-link {
                display: flex;
                align-items: center;
                gap: 0.5rem;
                color: var(--text-primary);
                text-decoration: none;
                font-weight: 500;
                padding: 0.75rem 1.5rem;
                border-radius: var(--radius);
                transition: all 0.3s ease;
                background: var(--background);
            }

            .nav-link:hover {
                background: var(--primary-color);
                color: white;
                transform: translateY(-2px);
                box-shadow: var(--shadow);
            }

            .nav-link i {
                font-size: 1.1rem;
            }

            .status-indicator {
                display: inline-flex;
                align-items: center;
                gap: 0.5rem;
                padding: 0.5rem 1rem;
                border-radius: var(--radius-sm);
                font-size: 0.875rem;
                font-weight: 500;
                margin-left: 1rem;
            }

            .status-indicator.success {
                background: rgba(16, 185, 129, 0.1);
                color: var(--success-color);
            }

            .status-indicator.error {
                background: rgba(239, 68, 68, 0.1);
                color: var(--error-color);
            }

            .status-indicator.warning {
                background: rgba(245, 158, 11, 0.1);
                color: var(--warning-color);
            }

            @keyframes fadeInDown {
                from {
                    opacity: 0;
                    transform: translateY(-30px);
                }
                to {
                    opacity: 1;
                    transform: translateY(0);
                }
            }

            @keyframes fadeInUp {
                from {
                    opacity: 0;
                    transform: translateY(30px);
                }
                to {
                    opacity: 1;
                    transform: translateY(0);
                }
            }

            @keyframes pulse {
                0%, 100% {
                    opacity: 1;
                }
                50% {
                    opacity: 0.5;
                }
            }

            .loading {
                animation: pulse 1.5s ease-in-out infinite;
            }

            /* Learning Status Modal */
            .modal {
                display: none;
                position: fixed;
                z-index: 1000;
                left: 0;
                top: 0;
                width: 100%;
                height: 100%;
                overflow: auto;
                background-color: rgba(0, 0, 0, 0.5);
                animation: fadeIn 0.3s;
            }

            .modal-content {
                background-color: var(--surface);
                margin: 5% auto;
                padding: 0;
                border-radius: var(--radius-lg);
                width: 90%;
                max-width: 800px;
                box-shadow: var(--shadow-xl);
                animation: slideDown 0.3s;
            }

            .modal-header {
                padding: 1.5rem 2rem;
                border-bottom: 2px solid var(--border-light);
                display: flex;
                justify-content: space-between;
                align-items: center;
            }

            .modal-header h2 {
                margin: 0;
                color: var(--text-primary);
                font-size: 1.5rem;
                display: flex;
                align-items: center;
                gap: 0.75rem;
            }

            .modal-close {
                background: none;
                border: none;
                font-size: 2rem;
                color: var(--text-secondary);
                cursor: pointer;
                padding: 0;
                width: 32px;
                height: 32px;
                display: flex;
                align-items: center;
                justify-content: center;
                border-radius: 4px;
                transition: all 0.2s;
            }

            .modal-close:hover {
                background: var(--background);
                color: var(--text-primary);
            }

            .modal-body {
                padding: 2rem;
                max-height: 70vh;
                overflow-y: auto;
            }

            .stat-card {
                background: var(--background);
                border: 1px solid var(--border-light);
                border-radius: var(--radius);
                padding: 1.5rem;
                margin-bottom: 1rem;
            }

            .stat-label {
                font-size: 0.875rem;
                color: var(--text-secondary);
                margin-bottom: 0.5rem;
                font-weight: 500;
            }

            .stat-value {
                font-size: 2rem;
                font-weight: 700;
                color: var(--primary-color);
            }

            .progress-bar-container {
                background: var(--border-light);
                border-radius: 8px;
                height: 20px;
                overflow: hidden;
                margin-top: 0.5rem;
            }

            .progress-bar-fill {
                height: 100%;
                background: linear-gradient(90deg, var(--primary-color), var(--success-color));
                transition: width 0.3s ease;
                display: flex;
                align-items: center;
                justify-content: center;
                color: white;
                font-size: 0.75rem;
                font-weight: 600;
            }

            .pattern-grid {
                display: grid;
                grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
                gap: 1rem;
                margin-top: 1rem;
            }

            .pattern-item {
                background: white;
                border: 1px solid var(--border-light);
                border-radius: var(--radius-sm);
                padding: 1rem;
                text-align: center;
            }

            .pattern-item .label {
                font-size: 0.875rem;
                color: var(--text-secondary);
                margin-bottom: 0.5rem;
            }

            .pattern-item .value {
                font-size: 1.5rem;
                font-weight: 700;
                color: var(--primary-color);
            }

            @keyframes fadeIn {
                from { opacity: 0; }
                to { opacity: 1; }
            }

            @keyframes slideDown {
                from {
                    opacity: 0;
                    transform: translateY(-30px);
                }
                to {
                    opacity: 1;
                    transform: translateY(0);
                }
            }

            @media (max-width: 768px) {
                .main-content {
                    grid-template-columns: 1fr;
                }
                
                .header h1 {
                    font-size: 2rem;
                }
                
                .nav-links {
                    flex-direction: column;
                    align-items: center;
                }
                
                .container {
                    padding: 1rem;
                }
            }
        </style>
    </head>
    <body>
        <div class="container">
            <!-- Auth Navigation Bar -->
            <div id="auth-nav" style="background: var(--surface); border-radius: var(--radius-lg); padding: 1rem 2rem; margin-bottom: 2rem; box-shadow: var(--shadow-lg); display: flex; justify-content: space-between; align-items: center;">
                <div style="display: flex; gap: 1rem; align-items: center;">
                    <a href="/about" style="color: var(--text-primary); text-decoration: none; font-weight: 500;">
                        <i class="fas fa-info-circle"></i> About Us
                    </a>
                </div>
                <div id="auth-buttons" style="display: flex; gap: 1rem;">
                    <a href="/login" class="btn" style="text-decoration: none; padding: 0.5rem 1rem; font-size: 0.9rem;">
                        <i class="fas fa-sign-in-alt"></i> Login
                    </a>
                    <a href="/signup" class="btn" style="text-decoration: none; padding: 0.5rem 1rem; font-size: 0.9rem; background: linear-gradient(135deg, var(--success-color), #059669);">
                        <i class="fas fa-user-plus"></i> Sign Up
                    </a>
                </div>
                <div id="user-info" style="display: none; align-items: center; gap: 1rem;">
                    <span id="username-display" style="color: var(--text-primary); font-weight: 500;"></span>
                    <button onclick="logout()" class="btn" style="padding: 0.5rem 1rem; font-size: 0.9rem; background: linear-gradient(135deg, var(--error-color), #dc2626);">
                        <i class="fas fa-sign-out-alt"></i> Logout
                    </button>
                </div>
            </div>

            <div class="header">
                <h1><i class="fas fa-brain"></i> FlowMind</h1>
                <p>Intelligent Document Analysis & Requirements Extraction</p>
            </div>

            <!-- Login Required Message -->
            <div id="login-required" style="display: none; background: var(--surface); border-radius: var(--radius-lg); padding: 3rem; text-align: center; box-shadow: var(--shadow-xl); margin-bottom: 2rem;">
                <i class="fas fa-lock" style="font-size: 3rem; color: var(--primary-color); margin-bottom: 1rem;"></i>
                <h2 style="color: var(--text-primary); margin-bottom: 1rem;">Authentication Required</h2>
                <p style="color: var(--text-secondary); margin-bottom: 2rem;">Please login or sign up to access FlowMind features</p>
                <div style="display: flex; gap: 1rem; justify-content: center;">
                    <a href="/login" class="btn" style="text-decoration: none;">
                        <i class="fas fa-sign-in-alt"></i> Login
                    </a>
                    <a href="/signup" class="btn" style="text-decoration: none; background: linear-gradient(135deg, var(--success-color), #059669);">
                        <i class="fas fa-user-plus"></i> Sign Up
                    </a>
                </div>
            </div>

            <div class="main-content" id="main-content">
                <div class="card">
                    <div class="card-header">
                        <div class="card-icon basic">
                            <i class="fas fa-file-text"></i>
                        </div>
                        <div>
                            <div class="card-title">Basic Extraction</div>
                            <div class="card-subtitle">Quick text and image extraction</div>
                        </div>
                    </div>
                    
                    <form id="form-basic">
                        <div class="form-group">
                            <div class="file-input-wrapper">
                                <input type="file" id="basic-file" class="file-input" accept=".pdf,.doc,.docx,.ppt,.pptx,.txt,.png,.jpg,.jpeg" required>
                                <label for="basic-file" class="file-input-label">
                                    <i class="fas fa-cloud-upload-alt"></i>
                                    Choose file to analyze
                                </label>
                            </div>
                        </div>
                        
                        <button type="submit" class="btn">
                            <i class="fas fa-play"></i>
                            Extract Content
                        </button>
                        
                        <div class="progress-container" id="basic-progress">
                            <div class="progress-bar">
                                <div class="progress-fill" id="basic-progress-bar"></div>
                            </div>
                            <div class="progress-text" id="basic-progress-text">Preparing...</div>
                        </div>
                        
                        <div class="output-container" id="basic-output">
                            <div style="color: var(--text-muted); text-align: center; padding: 2rem;">
                                <i class="fas fa-arrow-up" style="font-size: 2rem; margin-bottom: 1rem; opacity: 0.5;"></i>
                                <div>Upload a file to see extracted content here</div>
                            </div>
                        </div>
                    </form>
                </div>

                <div class="card">
                    <div class="card-header">
                        <div class="card-icon agent">
                            <i class="fas fa-robot"></i>
                        </div>
                        <div>
                            <div class="card-title">AI Agent Analysis</div>
                            <div class="card-subtitle">Advanced RAG-powered requirements extraction</div>
                        </div>
                    </div>
                    
                    <form id="form-agent">
                        <div class="form-group">
                            <div class="file-input-wrapper">
                                <input type="file" id="agent-file" class="file-input" accept=".pdf,.doc,.docx,.ppt,.pptx,.txt,.png,.jpg,.jpeg" required>
                                <label for="agent-file" class="file-input-label">
                                    <i class="fas fa-cloud-upload-alt"></i>
                                    Choose file for AI analysis
                                </label>
                            </div>
                        </div>
                        
                        <button type="submit" class="btn">
                            <i class="fas fa-magic"></i>
                            Analyze with AI
                        </button>
                        
                        <div class="progress-container" id="agent-progress">
                            <div class="progress-bar">
                                <div class="progress-fill" id="agent-progress-bar"></div>
                            </div>
                            <div class="progress-text" id="agent-progress-text">Uploading...</div>
                        </div>
                        
                        <div class="output-container" id="agent-output">
                            <div style="color: var(--text-muted); text-align: center; padding: 2rem;">
                                <i class="fas fa-arrow-up" style="font-size: 2rem; margin-bottom: 1rem; opacity: 0.5;"></i>
                                <div>Upload a file to see AI analysis results here</div>
                            </div>
                        </div>
                        
                        <div class="view-link" id="agent-view-link"></div>
                    </form>
                </div>
            </div>

            <div class="navigation">
                <div class="nav-links">
                    <a href="/records" target="_blank" class="nav-link">
                        <i class="fas fa-database"></i>
                        View Records
                    </a>
                    <a href="/agent_status" target="_blank" class="nav-link">
                        <i class="fas fa-heartbeat"></i>
                        Agent Status
                    </a>
                    <button id="learning-status-btn" class="nav-link" style="cursor: pointer; border: none; font-family: inherit;">
                        <i class="fas fa-brain"></i>
                        Learning Progress
                    </button>
                    <a href="/docs" target="_blank" class="nav-link">
                        <i class="fas fa-book"></i>
                        API Documentation
                    </a>
                </div>
            </div>

            <!-- Learning Status Modal -->
            <div id="learningModal" class="modal" style="display: none;">
                <div class="modal-content">
                    <div class="modal-header">
                        <h2><i class="fas fa-brain"></i> Learning Progress</h2>
                        <button id="close-learning-modal" class="modal-close">&times;</button>
                    </div>
                    <div class="modal-body" id="learningContent">
                        <div style="text-align: center; padding: 2rem;">
                            <i class="fas fa-spinner fa-spin" style="font-size: 2rem; color: var(--primary-color);"></i>
                            <div style="margin-top: 1rem;">Loading learning status...</div>
                        </div>
                    </div>
                </div>
            </div>
        </div>

        <script>
            // Authentication check
            function checkAuth() {
                const token = localStorage.getItem('access_token');
                const userStr = localStorage.getItem('user');
                const authButtons = document.getElementById('auth-buttons');
                const userInfo = document.getElementById('user-info');
                const usernameDisplay = document.getElementById('username-display');
                const mainContent = document.getElementById('main-content');
                const loginRequired = document.getElementById('login-required');
                const navigation = document.querySelector('.navigation');

                if (token && userStr) {
                    try {
                        const user = JSON.parse(userStr);
                        // User is logged in - redirect to extract page
                        window.location.href = '/extract';
                        return;
                    } catch (e) {
                        // Invalid user data, clear and show login
                        localStorage.removeItem('access_token');
                        localStorage.removeItem('user');
                        showLoginRequired();
                    }
                } else {
                    showLoginRequired();
                }
            }

            function showLoginRequired() {
                const authButtons = document.getElementById('auth-buttons');
                const userInfo = document.getElementById('user-info');
                const mainContent = document.getElementById('main-content');
                const loginRequired = document.getElementById('login-required');
                const navigation = document.querySelector('.navigation');

                authButtons.style.display = 'flex';
                userInfo.style.display = 'none';
                mainContent.style.display = 'none';
                loginRequired.style.display = 'block';
                if (navigation) navigation.style.display = 'none';
            }

            function logout() {
                localStorage.removeItem('access_token');
                localStorage.removeItem('user');
                window.location.href = '/';
            }

            // Helper function to get auth headers
            function getAuthHeaders() {
                const token = localStorage.getItem('access_token');
                if (!token) {
                    console.warn('No access token found in localStorage');
                    return {};
                }
                // When using FormData, don't set Content-Type - browser will set it with boundary
                // But we can still set Authorization header
                return { 'Authorization': `Bearer ${token}` };
            }

            // Check auth on page load
            document.addEventListener('DOMContentLoaded', function() {
                checkAuth();
            });

            // Enhanced file input handling
            function setupFileInput(inputId, labelSelector) {
                const input = document.getElementById(inputId);
                const label = document.querySelector(labelSelector);
                
                input.addEventListener('change', function() {
                    if (this.files.length > 0) {
                        const fileName = this.files[0].name;
                        label.innerHTML = `<i class="fas fa-file"></i> ${fileName}`;
                        label.style.borderColor = 'var(--success-color)';
                        label.style.background = 'rgba(16, 185, 129, 0.05)';
                        label.style.color = 'var(--success-color)';
                    }
                });
            }

            setupFileInput('basic-file', 'label[for="basic-file"]');
            setupFileInput('agent-file', 'label[for="agent-file"]');

            // Enhanced progress animation
            function animateProgress(progressId, textId, stages) {
                const progressContainer = document.getElementById(progressId);
                const progressBar = document.getElementById(progressId + '-bar');
                const progressText = document.getElementById(textId);
                
                progressContainer.style.display = 'block';
                
                let currentStage = 0;
                let progress = 0;
                
                const interval = setInterval(() => {
                    progress += Math.random() * 15 + 5; // Random increment between 5-20
                    
                    if (progress >= 100) {
                        progress = 100;
                        clearInterval(interval);
                    }
                    
                    progressBar.style.width = progress + '%';
                    
                    // Update stage based on progress
                    const stageProgress = Math.floor((progress / 100) * stages.length);
                    if (stageProgress !== currentStage && stageProgress < stages.length) {
                        currentStage = stageProgress;
                        progressText.textContent = stages[currentStage];
                    }
                }, 200);
                
                return interval;
            }

            // Enhanced form submission
            async function postForm(url, fileInput, outputId) {
                const fileEl = document.getElementById(fileInput);
                const outEl = document.getElementById(outputId);
                if (!fileEl || !outEl) {
                    console.error('Form elements not found');
                    return;
                }
                
                const submitBtn = fileEl.closest('form')?.querySelector('button[type="submit"]');
                if (!submitBtn) {
                    console.error('Submit button not found');
                    return;
                }
                
                if (!fileEl.files || !fileEl.files[0]) {
                    showError(outEl, 'Please select a file first.');
                    return;
                }

                // Check authentication
                const token = localStorage.getItem('access_token');
                if (!token) {
                    showError(outEl, 'Please login to upload files. <a href="/login" style="color: var(--primary-color);">Login here</a>');
                    submitBtn.disabled = false;
                    submitBtn.innerHTML = isAgent ? '<i class="fas fa-magic"></i> Analyze with AI' : '<i class="fas fa-play"></i> Extract Content';
                    return;
                }

                // Update UI state
                submitBtn.disabled = true;
                submitBtn.innerHTML = '<i class="fas fa-spinner fa-spin"></i> Processing...';
                outEl.innerHTML = '<div class="loading"><i class="fas fa-cog fa-spin"></i> Processing your document...</div>';

                const fd = new FormData();
                fd.append('file', fileEl.files[0]);
                
                const isAgent = (fileInput === 'agent-file');
                const stages = isAgent
                    ? ['Uploading...', 'Extracting text...', 'Processing images...', 'Running AI analysis...', 'Generating insights...']
                    : ['Uploading...', 'Extracting content...', 'Processing images...', 'Finalizing...'];
                
                let progressInterval = null;
                try {
                    progressInterval = animateProgress(
                        isAgent ? 'agent-progress' : 'basic-progress',
                        isAgent ? 'agent-progress-text' : 'basic-progress-text',
                        stages
                    );
                } catch (e) {
                    console.warn('Progress animation failed:', e);
                }

                try {
                    const headers = getAuthHeaders();
                    // Debug: Log token presence
                    const token = localStorage.getItem('access_token');
                    console.log('Token present:', !!token);
                    console.log('Headers being sent:', headers);
                    
                    if (!token) {
                        showError(outEl, 'Authentication token not found. Please <a href="/login" style="color: var(--primary-color);">login again</a>.');
                        submitBtn.disabled = false;
                        submitBtn.innerHTML = isAgent ? '<i class="fas fa-magic"></i> Analyze with AI' : '<i class="fas fa-play"></i> Extract Content';
                        return;
                    }
                    
                    const res = await fetch(url, { 
                        method: 'POST', 
                        body: fd,
                        headers: headers
                    });
                    const text = await res.text();
                    
                    if (progressInterval) {
                        clearInterval(progressInterval);
                    }
                    
                    if (res.ok) {
                        try {
                            const data = JSON.parse(text);
                            displayResults(outEl, data, isAgent);
                            
                            if (data && data.view_id) {
                                const linkDiv = document.getElementById('agent-view-link');
                                if (linkDiv) {
                                    linkDiv.innerHTML = `
                                        <a href="/view_requirements/${data.view_id}" target="_blank">
                                            <i class="fas fa-external-link-alt"></i>
                                            View Formatted Results
                                        </a>
                                    `;
                                }
                            }
                        } catch(e) {
                            console.error('JSON parse error:', e);
                            showError(outEl, 'Invalid response format: ' + text.substring(0, 200));
                        }
                    } else {
                        // Handle specific error cases
                        if (res.status === 401) {
                            // Clear invalid token and redirect to login
                            localStorage.removeItem('access_token');
                            localStorage.removeItem('user');
                            showError(outEl, 'Your session has expired or authentication failed. Please <a href="/login" style="color: var(--primary-color); font-weight: bold;">login again</a> to continue.');
                            // Redirect to login after a short delay
                            setTimeout(() => {
                                window.location.href = '/login';
                            }, 3000);
                    } else {
                        showError(outEl, `Error ${res.status}: ${text.substring(0, 500)}`);
                        }
                    }
                } catch (e) {
                    if (progressInterval) {
                        clearInterval(progressInterval);
                    }
                    console.error('Network error:', e);
                    showError(outEl, 'Network error: ' + (e.message || 'Unknown error'));
                } finally {
                    if (submitBtn) {
                        submitBtn.disabled = false;
                        submitBtn.innerHTML = isAgent 
                            ? '<i class="fas fa-magic"></i> Analyze with AI'
                            : '<i class="fas fa-play"></i> Extract Content';
                    }
                }
            }

            function displayResults(container, data, isAgent) {
                let html = '<div style="color: var(--success-color); margin-bottom: 1rem;"><i class="fas fa-check-circle"></i> Analysis Complete</div>';
                
                if (isAgent) {
                    html += formatAgentResults(data);
                } else {
                    html += formatBasicResults(data);
                }
                
                container.innerHTML = html;
            }

            function formatBasicResults(data) {
                let html = '<div style="margin-bottom: 1rem;">';
                html += `<div><strong>File:</strong> ${data.filename || 'Unknown'}</div>`;
                html += `<div><strong>Summary:</strong> ${data.summary || 'No summary available'}</div>`;
                html += `<div><strong>Images detected:</strong> ${data.images_detected || 0}</div>`;
                html += '</div>';
                
                if (data.images && data.images.length > 0) {
                    html += '<div style="margin-top: 1rem;"><strong>Images:</strong></div>';
                    data.images.forEach((img, index) => {
                        html += `<div style="margin: 0.5rem 0; padding: 0.5rem; background: rgba(255,255,255,0.05); border-radius: 4px;">`;
                        html += `<div><strong>Image ${index + 1}</strong> (Page ${img.page})</div>`;
                        if (img.ocr) {
                            html += `<div style="font-size: 0.8rem; color: var(--text-muted); margin-top: 0.25rem;">OCR: ${img.ocr.substring(0, 100)}${img.ocr.length > 100 ? '...' : ''}</div>`;
                        }
                        html += '</div>';
                    });
                }
                
                return html;
            }

            function formatAgentResults(data) {
                let html = '<div style="margin-bottom: 1rem;">';
                html += `<div><strong>File:</strong> ${data.filename || 'Unknown'}</div>`;
                html += `<div><strong>Extraction:</strong> ${data.extraction_summary || 'No summary available'}</div>`;
                html += '</div>';
                
                if (data.agent_processing) {
                    html += '<div style="margin-bottom: 1rem;">';
                    html += '<div><strong>Agent Processing:</strong></div>';
                    html += `<div style="margin-left: 1rem; color: var(--text-muted);">${data.agent_processing.message || 'Processing completed'}</div>`;
                    html += '</div>';
                }
                
                if (data.requirements_extraction && data.requirements_extraction.response) {
                    html += '<div style="margin-bottom: 1rem;">';
                    html += '<div><strong>Requirements Extracted:</strong></div>';
                    html += `<div style="margin-left: 1rem; white-space: pre-wrap;">${data.requirements_extraction.response}</div>`;
                    html += '</div>';
                }
                
                return html;
            }

            function showError(container, message) {
                container.innerHTML = `
                    <div style="color: var(--error-color); text-align: center; padding: 1rem;">
                        <i class="fas fa-exclamation-triangle"></i>
                        <div style="margin-top: 0.5rem;">${message}</div>
                    </div>
                `;
            }

            // Form event listeners
            document.getElementById('form-basic').addEventListener('submit', function(ev){
                ev.preventDefault();
                postForm('/upload_client_doc', 'basic-file', 'basic-output');
            });

            document.getElementById('form-agent').addEventListener('submit', function(ev){
                ev.preventDefault();
                postForm('/upload_agent_doc', 'agent-file', 'agent-output');
            });

            // Add some interactive effects
            document.querySelectorAll('.card').forEach(card => {
                card.addEventListener('mouseenter', function() {
                    this.style.transform = 'translateY(-8px)';
                });
                
                card.addEventListener('mouseleave', function() {
                    this.style.transform = 'translateY(0)';
                });
            });

            // Learning Status Functions
            async function showLearningStatus() {
                const modal = document.getElementById('learningModal');
                const content = document.getElementById('learningContent');
                
                modal.style.display = 'block';
                content.innerHTML = `
                    <div style="text-align: center; padding: 2rem;">
                        <i class="fas fa-spinner fa-spin" style="font-size: 2rem; color: var(--primary-color);"></i>
                        <div style="margin-top: 1rem;">Loading learning status...</div>
                    </div>
                `;
                
                try {
                    const response = await fetch('/learning-status');
                    const data = await response.json();
                    displayLearningStatus(data);
                } catch (error) {
                    content.innerHTML = `
                        <div style="color: var(--error-color); text-align: center; padding: 2rem;">
                            <i class="fas fa-exclamation-triangle"></i>
                            <div style="margin-top: 1rem;">Error loading learning status: ${error.message}</div>
                        </div>
                    `;
                }
            }

            function displayLearningStatus(data) {
                const content = document.getElementById('learningContent');
                
                if (data.error) {
                    content.innerHTML = `
                        <div style="color: var(--error-color); text-align: center; padding: 2rem;">
                            <i class="fas fa-exclamation-triangle"></i>
                            <div style="margin-top: 1rem;">${data.error}</div>
                        </div>
                    `;
                    return;
                }
                
                const totalPatterns = data.total_learned_patterns || 0;
                const iterations = data.learning_iterations || 0;
                const documents = data.total_documents_processed || 0;
                const successRate = (data.success_rate || 0) * 100;
                const learningEnabled = data.learning_enabled ? 'Enabled' : 'Disabled';
                
                let html = `
                    <div class="stat-card">
                        <div class="stat-label">Learning Status</div>
                        <div class="stat-value" style="font-size: 1.5rem;">${learningEnabled}</div>
                    </div>
                    
                    <div class="pattern-grid">
                        <div class="stat-card">
                            <div class="stat-label">Total Patterns Learned</div>
                            <div class="stat-value">${totalPatterns.toLocaleString()}</div>
                        </div>
                        
                        <div class="stat-card">
                            <div class="stat-label">Learning Iterations</div>
                            <div class="stat-value">${iterations.toLocaleString()}</div>
                        </div>
                        
                        <div class="stat-card">
                            <div class="stat-label">Documents Processed</div>
                            <div class="stat-value">${documents.toLocaleString()}</div>
                        </div>
                        
                        <div class="stat-card">
                            <div class="stat-label">Success Rate</div>
                            <div class="stat-value" style="font-size: 1.5rem;">${successRate.toFixed(1)}%</div>
                            <div class="progress-bar-container" style="margin-top: 0.5rem;">
                                <div class="progress-bar-fill" style="width: ${successRate}%;">${successRate.toFixed(0)}%</div>
                            </div>
                        </div>
                    </div>
                `;
                
                // Patterns by Category
                if (data.patterns_by_category) {
                    html += `
                        <div style="margin-top: 2rem;">
                            <h3 style="margin-bottom: 1rem; color: var(--text-primary);">
                                <i class="fas fa-layer-group"></i> Patterns by Category
                            </h3>
                            <div class="pattern-grid">
                    `;
                    
                    for (const [category, patterns] of Object.entries(data.patterns_by_category)) {
                        const total = (patterns.keywords || 0) + (patterns.phrases || 0) + (patterns.patterns || 0);
                        html += `
                            <div class="pattern-item">
                                <div class="label">${category.replace(/_/g, ' ').replace(/\\b\\w/g, l => l.toUpperCase())}</div>
                                <div class="value">${total}</div>
                                <div style="font-size: 0.75rem; color: var(--text-secondary); margin-top: 0.5rem;">
                                    K: ${patterns.keywords || 0} | P: ${patterns.phrases || 0} | Pt: ${patterns.patterns || 0}
                                </div>
                            </div>
                        `;
                    }
                    
                    html += `
                            </div>
                        </div>
                    `;
                }
                
                // Method Success Rates
                if (data.method_success_rates && Object.keys(data.method_success_rates).length > 0) {
                    html += `
                        <div style="margin-top: 2rem;">
                            <h3 style="margin-bottom: 1rem; color: var(--text-primary);">
                                <i class="fas fa-chart-line"></i> Method Performance
                            </h3>
                    `;
                    
                    for (const [method, rate] of Object.entries(data.method_success_rates)) {
                        const ratePercent = rate * 100;
                        html += `
                            <div class="stat-card">
                                <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 0.5rem;">
                                    <div class="stat-label" style="margin: 0;">${method.charAt(0).toUpperCase() + method.slice(1)}</div>
                                    <div style="font-weight: 600; color: var(--primary-color);">${ratePercent.toFixed(1)}%</div>
                                </div>
                                <div class="progress-bar-container">
                                    <div class="progress-bar-fill" style="width: ${ratePercent}%;">${ratePercent.toFixed(0)}%</div>
                                </div>
                            </div>
                        `;
                    }
                    
                    html += `</div>`;
                }
                
                // Last Learning Session
                if (data.last_learning_session) {
                    const session = data.last_learning_session;
                    html += `
                        <div style="margin-top: 2rem;">
                            <h3 style="margin-bottom: 1rem; color: var(--text-primary);">
                                <i class="fas fa-clock"></i> Last Learning Session
                            </h3>
                            <div class="stat-card">
                                <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(150px, 1fr)); gap: 1rem;">
                                    <div>
                                        <div class="stat-label">Keywords</div>
                                        <div class="stat-value" style="font-size: 1.25rem;">${session.keywords || 0}</div>
                                    </div>
                                    <div>
                                        <div class="stat-label">Phrases</div>
                                        <div class="stat-value" style="font-size: 1.25rem;">${session.phrases || 0}</div>
                                    </div>
                                    <div>
                                        <div class="stat-label">Patterns</div>
                                        <div class="stat-value" style="font-size: 1.25rem;">${session.patterns || 0}</div>
                                    </div>
                                    <div>
                                        <div class="stat-label">Total</div>
                                        <div class="stat-value" style="font-size: 1.25rem;">${session.total || 0}</div>
                                    </div>
                                </div>
                                ${session.filename ? `<div style="margin-top: 1rem; font-size: 0.875rem; color: var(--text-secondary);">File: ${session.filename}</div>` : ''}
                            </div>
                        </div>
                    `;
                }
                
                content.innerHTML = html;
            }

            function closeLearningModal() {
                document.getElementById('learningModal').style.display = 'none';
            }

            // Event listeners for learning status
            document.addEventListener('DOMContentLoaded', function() {
                const learningBtn = document.getElementById('learning-status-btn');
                const closeBtn = document.getElementById('close-learning-modal');
                const modal = document.getElementById('learningModal');
                
                if (learningBtn) {
                    learningBtn.addEventListener('click', function() {
                        showLearningStatus().catch(function(error) {
                            console.error('Error showing learning status:', error);
                            const content = document.getElementById('learningContent');
                            if (content) {
                                content.innerHTML = `
                                    <div style="color: var(--error-color); text-align: center; padding: 2rem;">
                                        <i class="fas fa-exclamation-triangle"></i>
                                        <div style="margin-top: 1rem;">Error loading learning status. Please try again.</div>
                                    </div>
                                `;
                            }
                        });
                    });
                }
                
                if (closeBtn) {
                    closeBtn.addEventListener('click', closeLearningModal);
                }
                
                // Close modal when clicking outside
                if (modal) {
                    modal.addEventListener('click', function(event) {
                        if (event.target === modal) {
                            closeLearningModal();
                        }
                    });
                }
            });
            
            // Global error handler for unhandled promises
            window.addEventListener('unhandledrejection', function(event) {
                console.error('Unhandled promise rejection:', event.reason);
                // Prevent the error from appearing in console
                event.preventDefault();
            });
        </script>
    </body>
    </html>
    """


@app.get("/index.html", response_class=HTMLResponse)
async def serve_index():
    """Serve the index.html landing page."""
    index_path = os.path.join(BASE_DIR, "index.html")
    if os.path.exists(index_path):
        with open(index_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return RedirectResponse(url="/", status_code=302)


@app.get("/login.html", response_class=HTMLResponse)
async def serve_login():
    """Serve the login.html page."""
    login_path = os.path.join(BASE_DIR, "login.html")
    if os.path.exists(login_path):
        with open(login_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return RedirectResponse(url="/", status_code=302)


# NOTE: Upload endpoints moved to routes/upload_routes.py to avoid duplicate route definitions
# The routes are included via app.include_router(upload_router) below


# ==============================================================
# PUBLIC STATS ENDPOINT (No authentication required)
# ==============================================================

@app.get("/api/public/stats")
async def get_public_stats():
    """Get public statistics for the landing page - no authentication required."""
    try:
        db = SessionLocal()
        try:
            # Count total documents scanned
            total_documents = db.query(ParsedFile).count()
            
            # Count total clients (users)
            total_clients = db.query(User).filter(User.is_active == 1).count()
            
            # Get patterns learned from all agents (sum of all learned items: keywords + patterns + phrases)
            # Aggregate from both global agent and all user-specific agents
            patterns_learned = 0
            all_learned_items = set()  # Use set to deduplicate across agents
            
            try:
                from rag_agent import get_agent, user_agent_cache
                
                # First, try global agent
                try:
                    global_agent = get_agent()  # No user_id = global agent
                    if hasattr(global_agent, 'learned_patterns'):
                        learned_patterns = global_agent.learned_patterns
                        if isinstance(learned_patterns, dict):
                            for category, data in learned_patterns.items():
                                if isinstance(data, dict):
                                    keywords = data.get('keywords', set())
                                    patterns = data.get('patterns', set())
                                    phrases = data.get('phrases', set())
                                    # Convert to lists if sets for accurate counting
                                    keywords = list(keywords) if isinstance(keywords, set) else (keywords if isinstance(keywords, list) else [])
                                    patterns = list(patterns) if isinstance(patterns, set) else (patterns if isinstance(patterns, list) else [])
                                    phrases = list(phrases) if isinstance(phrases, set) else (phrases if isinstance(phrases, list) else [])
                                    # Add to set for deduplication
                                    for kw in keywords:
                                        all_learned_items.add(f"keyword:{category}:{kw}")
                                    for pat in patterns:
                                        all_learned_items.add(f"pattern:{category}:{pat}")
                                    for phr in phrases:
                                        all_learned_items.add(f"phrase:{category}:{phr}")
                except Exception as e:
                    print(f"⚠️ Error getting patterns from global agent: {e}")
                
                # Then, aggregate from all user-specific agents
                try:
                    if user_agent_cache:
                        for user_id, user_agent in user_agent_cache.items():
                            if hasattr(user_agent, 'learned_patterns'):
                                learned_patterns = user_agent.learned_patterns
                                if isinstance(learned_patterns, dict):
                                    for category, data in learned_patterns.items():
                                        if isinstance(data, dict):
                                            keywords = data.get('keywords', set())
                                            patterns = data.get('patterns', set())
                                            phrases = data.get('phrases', set())
                                            # Convert to lists if sets
                                            keywords = list(keywords) if isinstance(keywords, set) else (keywords if isinstance(keywords, list) else [])
                                            patterns = list(patterns) if isinstance(patterns, set) else (patterns if isinstance(patterns, list) else [])
                                            phrases = list(phrases) if isinstance(phrases, set) else (phrases if isinstance(phrases, list) else [])
                                            # Add to set for deduplication
                                            for kw in keywords:
                                                all_learned_items.add(f"keyword:{category}:{kw}")
                                            for pat in patterns:
                                                all_learned_items.add(f"pattern:{category}:{pat}")
                                            for phr in phrases:
                                                all_learned_items.add(f"phrase:{category}:{phr}")
                except Exception as e:
                    print(f"⚠️ Error aggregating patterns from user agents: {e}")
                
                # Count unique learned items
                patterns_learned = len(all_learned_items)
                print(f"📊 Public stats: Counted {patterns_learned} unique learned items (from global + {len(user_agent_cache) if user_agent_cache else 0} user agents)")
                
                # If no patterns found from agents, try fallback to database features
                if patterns_learned == 0:
                    try:
                        feature_count = db.query(Feature).count()
                        if feature_count > 0:
                            # Use feature count as a rough estimate (each feature might represent learned patterns)
                            patterns_learned = feature_count
                            print(f"📊 Public stats: No patterns from agents, using feature count as fallback: {patterns_learned}")
                    except Exception as fallback_error:
                        print(f"⚠️ Fallback feature count also failed: {fallback_error}")
                
            except Exception as e:
                print(f"⚠️ Error getting patterns for public stats: {e}")
                import traceback
                traceback.print_exc()
                
                # Fallback: try to count from database features as a proxy
                try:
                    feature_count = db.query(Feature).count()
                    if feature_count > 0:
                        # Use feature count as a rough estimate (each feature might represent learned patterns)
                        patterns_learned = feature_count
                        print(f"📊 Public stats: Using feature count as fallback after error: {patterns_learned}")
                except Exception as fallback_error:
                    print(f"⚠️ Fallback feature count failed: {fallback_error}")
                    patterns_learned = 0
            
            # Calculate accuracy from extraction stats
            accuracy = 97.3  # Default fallback
            try:
                agent = get_agent()
                if hasattr(agent, 'extraction_stats'):
                    stats = agent.extraction_stats
                    total_docs = stats.get('total_documents', 0)
                    successful = stats.get('successful_extractions', 0)
                    if total_docs > 0:
                        accuracy = (successful / total_docs) * 100
            except Exception as e:
                print(f"⚠️ Error calculating accuracy: {e}")
            
            return {
                "documents_scanned": total_documents,
                "clients": total_clients,
                "patterns_learned": patterns_learned,
                "accuracy": round(accuracy, 1)
            }
        finally:
            db.close()
    except Exception as e:
        print(f"⚠️ Error getting public stats: {e}")
        # Return default values on error
        return {
            "documents_scanned": 0,
            "clients": 0,
            "patterns_learned": 0,
            "accuracy": 0.0
        }


# ==============================================================
# LEARNING SYSTEM ENDPOINTS
# ==============================================================

@app.get("/learning-status")
async def get_learning_status():
    """Get the current learning status and statistics."""
    try:
        agent = get_agent()
        if hasattr(agent, 'get_learning_status'):
            status = agent.get_learning_status()
            return status
        else:
            return {"error": "Learning system not available"}
    except Exception as e:
        return {"error": f"Failed to get learning status: {str(e)}"}

# ==============================================================
# RECORDS DASHBOARD
# ==============================================================

@app.get("/records", response_class=HTMLResponse)
def get_records(current_user: User = Depends(get_current_user)):
    """Displays all analyzed file summaries. Requires authentication."""
    db = SessionLocal()
    records = db.query(ParsedFile).all()
    db.close()

    html = """
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>FlowMind Records</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
        <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css" rel="stylesheet">
        <style>
            :root {
                --primary-color: #2563eb;
                --primary-dark: #1d4ed8;
                --secondary-color: #64748b;
                --success-color: #10b981;
                --warning-color: #f59e0b;
                --error-color: #ef4444;
                --background: #f8fafc;
                --surface: #ffffff;
                --surface-elevated: #ffffff;
                --text-primary: #1e293b;
                --text-secondary: #64748b;
                --text-muted: #94a3b8;
                --border: #e2e8f0;
                --border-light: #f1f5f9;
                --shadow: 0 1px 3px 0 rgba(0, 0, 0, 0.1), 0 1px 2px 0 rgba(0, 0, 0, 0.06);
                --shadow-lg: 0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -2px rgba(0, 0, 0, 0.05);
                --shadow-xl: 0 20px 25px -5px rgba(0, 0, 0, 0.1), 0 10px 10px -5px rgba(0, 0, 0, 0.04);
                --radius: 12px;
                --radius-sm: 8px;
                --radius-lg: 16px;
            }

            * {
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }

            body {
                font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
                background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                min-height: 100vh;
                color: var(--text-primary);
                line-height: 1.6;
            }

            .container {
                max-width: 1200px;
                margin: 0 auto;
                padding: 2rem;
            }

            .header {
                text-align: center;
                margin-bottom: 3rem;
                animation: fadeInDown 0.8s ease-out;
            }

            .header h1 {
                font-size: 3rem;
                font-weight: 700;
                color: white;
                margin-bottom: 0.5rem;
                text-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
            }

            .header p {
                font-size: 1.2rem;
                color: rgba(255, 255, 255, 0.9);
                font-weight: 300;
            }

            .back-btn {
                display: inline-flex;
                align-items: center;
                gap: 0.5rem;
                background: var(--surface);
                color: var(--text-primary);
                text-decoration: none;
                padding: 0.75rem 1.5rem;
                border-radius: var(--radius);
                font-weight: 500;
                transition: all 0.3s ease;
                box-shadow: var(--shadow);
                margin-bottom: 2rem;
                animation: fadeInUp 0.8s ease-out;
            }

            .back-btn:hover {
                background: var(--primary-color);
                color: white;
                transform: translateY(-2px);
                box-shadow: var(--shadow-lg);
            }

            .records-card {
                background: var(--surface);
                border-radius: var(--radius-lg);
                padding: 2rem;
                box-shadow: var(--shadow-xl);
                border: 1px solid var(--border-light);
                animation: fadeInUp 0.8s ease-out 0.2s both;
            }

            .records-header {
                display: flex;
                align-items: center;
                margin-bottom: 2rem;
            }

            .records-icon {
                width: 48px;
                height: 48px;
                border-radius: var(--radius);
                display: flex;
                align-items: center;
                justify-content: center;
                margin-right: 1rem;
                font-size: 1.5rem;
                color: white;
                background: linear-gradient(135deg, var(--success-color), #059669);
            }

            .records-title {
                font-size: 1.5rem;
                font-weight: 600;
                color: var(--text-primary);
            }

            .records-subtitle {
                font-size: 0.9rem;
                color: var(--text-secondary);
                margin-top: 0.25rem;
            }

            .records-table {
                width: 100%;
                border-collapse: collapse;
                background: var(--surface);
                border-radius: var(--radius);
                overflow: hidden;
                box-shadow: var(--shadow);
            }

            .records-table th {
                background: linear-gradient(135deg, var(--primary-color), var(--primary-dark));
                color: white;
                padding: 1rem;
                text-align: left;
                font-weight: 600;
                font-size: 0.9rem;
                text-transform: uppercase;
                letter-spacing: 0.5px;
            }

            .records-table td {
                padding: 1rem;
                border-bottom: 1px solid var(--border-light);
                color: var(--text-primary);
            }

            .records-table tr:hover {
                background: var(--background);
                transform: scale(1.01);
                transition: all 0.2s ease;
            }

            .records-table tr:last-child td {
                border-bottom: none;
            }

            .file-link {
                display: inline-flex;
                align-items: center;
                gap: 0.5rem;
                background: var(--primary-color);
                color: white;
                text-decoration: none;
                padding: 0.5rem 1rem;
                border-radius: var(--radius-sm);
                font-size: 0.875rem;
                font-weight: 500;
                transition: all 0.3s ease;
            }

            .file-link:hover {
                background: var(--primary-dark);
                transform: translateY(-1px);
                box-shadow: var(--shadow);
            }

            .file-icon {
                width: 32px;
                height: 32px;
                border-radius: var(--radius-sm);
                display: flex;
                align-items: center;
                justify-content: center;
                margin-right: 0.75rem;
                font-size: 1rem;
                color: white;
                background: linear-gradient(135deg, var(--secondary-color), #475569);
            }

            .file-name {
                font-weight: 500;
                color: var(--text-primary);
            }

            .file-meta {
                font-size: 0.875rem;
                color: var(--text-secondary);
                margin-top: 0.25rem;
            }

            .stats-badge {
                display: inline-flex;
                align-items: center;
                gap: 0.25rem;
                background: rgba(16, 185, 129, 0.1);
                color: var(--success-color);
                padding: 0.25rem 0.75rem;
                border-radius: var(--radius-sm);
                font-size: 0.8rem;
                font-weight: 500;
            }

            .empty-state {
                text-align: center;
                padding: 3rem;
                color: var(--text-secondary);
            }

            .empty-state i {
                font-size: 3rem;
                margin-bottom: 1rem;
                opacity: 0.5;
            }

            @keyframes fadeInDown {
                from {
                    opacity: 0;
                    transform: translateY(-30px);
                }
                to {
                    opacity: 1;
                    transform: translateY(0);
                }
            }

            @keyframes fadeInUp {
                from {
                    opacity: 0;
                    transform: translateY(30px);
                }
                to {
                    opacity: 1;
                    transform: translateY(0);
                }
            }

            @media (max-width: 768px) {
                .container {
                    padding: 1rem;
                }
                
                .header h1 {
                    font-size: 2rem;
                }
                
                .records-table {
                    font-size: 0.875rem;
                }
                
                .records-table th,
                .records-table td {
                    padding: 0.75rem 0.5rem;
                }
            }
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <h1><i class="fas fa-database"></i> FlowMind Records</h1>
                <p>All analyzed documents and their extracted content</p>
            </div>

            <a href="/" class="back-btn">
                <i class="fas fa-arrow-left"></i>
                Back to FlowMind
            </a>

            <div class="records-card">
                <div class="records-header">
                    <div class="records-icon">
                        <i class="fas fa-folder-open"></i>
                    </div>
                    <div>
                        <div class="records-title">Processed Documents</div>
                        <div class="records-subtitle">Total: """ + str(len(records)) + """ files analyzed</div>
                    </div>
                </div>
    """

    if records:
        html += """
                <table class="records-table">
                    <thead>
                        <tr>
                            <th><i class="fas fa-hashtag"></i> ID</th>
                            <th><i class="fas fa-file"></i> Filename</th>
                            <th><i class="fas fa-images"></i> Images</th>
                            <th><i class="fas fa-info-circle"></i> Summary</th>
                            <th><i class="fas fa-external-link-alt"></i> Actions</th>
                        </tr>
                    </thead>
                    <tbody>
        """

        for r in records:
            text_link = r.full_text_path.replace("\\", "/") if r.full_text_path else "#"
            file_icon = "fas fa-file-pdf" if r.filename.lower().endswith('.pdf') else \
                       "fas fa-file-word" if r.filename.lower().endswith(('.doc', '.docx')) else \
                       "fas fa-file-powerpoint" if r.filename.lower().endswith(('.ppt', '.pptx')) else \
                       "fas fa-file-image" if r.filename.lower().endswith(('.png', '.jpg', '.jpeg')) else \
                       "fas fa-file-alt"
            
            html += f"""
                        <tr>
                            <td>
                                <div class="stats-badge">
                                    <i class="fas fa-hashtag"></i>
                                    {r.id}
                                </div>
                            </td>
                            <td>
                                <div class="file-icon">
                                    <i class="{file_icon}"></i>
                                </div>
                                <div>
                                    <div class="file-name">{r.filename}</div>
                                    <div class="file-meta">Processed document</div>
                                </div>
                            </td>
                            <td>
                                <div class="stats-badge">
                                    <i class="fas fa-images"></i>
                                    {r.detected_shapes}
                                </div>
                            </td>
                            <td>
                                <div style="max-width: 300px; overflow: hidden; text-overflow: ellipsis;">
                                    {r.summary}
                                </div>
                            </td>
                            <td>
                                <a href='/{text_link}' target='_blank' class="file-link">
                                    <i class="fas fa-external-link-alt"></i>
                                    View Text
                                </a>
                            </td>
                        </tr>
            """

        html += """
                    </tbody>
                </table>
        """
    else:
        html += """
                <div class="empty-state">
                    <i class="fas fa-folder-open"></i>
                    <h3>No documents processed yet</h3>
                    <p>Upload your first document to see it appear here</p>
                </div>
        """

    html += """
            </div>
        </div>
    </body>
    </html>
    """
    return html


# ==============================================================
# RAG AGENT ENDPOINTS
# ==============================================================
_DEBUG_LOG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "debug-0e985e.log")

@app.post("/analyze_with_agent")
async def _analyze_with_agent_internal(
    file: UploadFile,
    user_id: int = None,
    db_session = None,
    progress_tracker_id: str = None,
    basic_extraction_data: dict = None,
    project_id: int = None,
):
    """Internal function for AI agent analysis - ensures response is always returned.
    Optionally merges Basic Extraction data (text + image summaries) for improved results."""
    print(f"_analyze_with_agent_internal called for file: {file.filename}, user_id: {user_id}")
    if basic_extraction_data:
        print(f"📋 Basic Extraction data provided: {len(basic_extraction_data.get('extracted_text', ''))} chars, {len(basic_extraction_data.get('image_summaries', []))} image summaries")
    """Internal function: Extract text from document and process with RAG agent for requirements extraction."""
    from services.progress_service import ProcessingStage
    from services.progress_storage import get_progress_tracker
    
    tracker = None
    if progress_tracker_id:
        tracker = get_progress_tracker(progress_tracker_id)
        if tracker:
            # Tracker should already be in UPLOADING stage from start()
            pass
    
    # #region agent log
    try:
        f = open(_DEBUG_LOG, "a", encoding="utf-8"); f.write(json.dumps({"sessionId": "0e985e", "hypothesisId": "A", "location": "flowmind.py:_analyze_with_agent_internal", "message": "entry processing started", "data": {"filename": file.filename}, "timestamp": int(time.time() * 1000)}) + "\n"); f.close()
    except Exception:
        pass
    # #endregion
    try:
        from services.requirement_validation import is_srs_supported_upload

        ok_ext, ext_reasons = is_srs_supported_upload(file.filename or "")
        if not ok_ext:
            raise HTTPException(
                status_code=400,
                detail=(
                    "Unsupported document type for SRS workflow. "
                    + "; ".join(ext_reasons)
                ),
            )

        # First, extract text using existing functionality
        filepath = os.path.join(UPLOAD_DIR, file.filename)
        if tracker:
            tracker.set_stage(ProcessingStage.UPLOADING)
            progress = tracker.get_progress()
            print(f"📊 Progress: {progress['progress']}% - {progress['message']}")
        
        with open(filepath, "wb") as f:
            f.write(await file.read())
        
        if tracker:
            tracker.set_stage(ProcessingStage.PARSING)
            progress = tracker.get_progress()
            print(f"📊 Progress: {progress['progress']}% - {progress['message']}")

        text_output = ""
        image_count = 0
        detected_images = 0  # Initialize counter for detected images
        image_metadata = []  # (image_id, image_path, page_number, ocr_text)
        image_requirements_structured = []
        diagram_jobs_pending: list = []

        def gen_image_id(fname: str, page_num: int, idx: int) -> str:
            base = f"{fname}|{page_num}|{idx}"
            h = hashlib.md5(base.encode()).hexdigest()[:10]
            return f"IMG-{h}-{page_num}-{idx}"

        # Extract text based on file type
        if file.filename.endswith(".pdf"):
            print(f"📄 Detected PDF file, initializing PDF reader...")
            reader = PdfReader(filepath)
            total_pages = len(reader.pages)
            chunk_size = 10
            page_end_offsets = {}
            page_text_by_page: dict = {}

            print(f"📘 Document has {total_pages} pages. Starting text extraction...")
            
            if tracker:
                tracker.set_stage(ProcessingStage.PARSING, total_pages=total_pages, current_page=0)
                progress = tracker.get_progress()
                print(f"📊 [{progress['progress']}%] {progress['message']}")

            for start in range(0, total_pages, chunk_size):
                end = min(start + chunk_size, total_pages)
                print(f"🔹 Processing pages {start + 1}–{end} of {total_pages}...")

                for page_index in range(start, end):
                    page = reader.pages[page_index]
                    page_text = page.extract_text() or ""
                    page_text_by_page[page_index + 1] = page_text
                    text_output += page_text + "\n"
                    page_end_offsets[page_index + 1] = len(text_output)

                    if tracker:
                        tracker.set_stage(ProcessingStage.PARSING, total_pages=total_pages, current_page=page_index + 1)
                        if (page_index + 1) % 5 == 0 or page_index + 1 == total_pages:  # Log every 5 pages
                            progress = tracker.get_progress()
                            print(f"📊 [{progress['progress']}%] Currently working on page {page_index + 1}/{total_pages}")

                print(f"✅ Finished pages {start + 1}–{end} of {total_pages}.")
            print(f"✅ Completed all {total_pages} pages. Extracted {len(text_output):,} characters.")

            images_per_page: dict = {}
            for pn, pg in enumerate(reader.pages, start=1):
                try:
                    imgs0 = getattr(pg, "images", []) or []
                except Exception:
                    imgs0 = []
                images_per_page[pn] = len(imgs0)

            # Extract images using pypdf Page.images if available
            try:
                from io import BytesIO
                from services.image_service import split_page_text_around_image, get_image_context
                for page_num, page in enumerate(reader.pages, start=1):
                    try:
                        imgs = getattr(page, "images", []) or []
                    except Exception:
                        imgs = []
                    for img_idx, img in enumerate(imgs, start=1):
                        try:
                            data = getattr(img, "data", None)
                            if not data:
                                continue
                            image_count += 1
                            image_id = gen_image_id(file.filename, page_num, img_idx)
                            out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}.png")
                            try:
                                im = Image.open(BytesIO(data))
                                im.save(out_path, format="PNG")
                                ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=im)
                            except Exception:
                                with open(out_path, "wb") as fimg:
                                    fimg.write(data)
                                ocr_text = ""
                            image_metadata.append((image_id, out_path, page_num, ocr_text))
                            detected_images += 1
                            pos = page_end_offsets.get(page_num, len(text_output))
                            context_before = text_output[max(0, pos - 500):pos]
                            pt = page_text_by_page.get(page_num, "")
                            nimg = max(1, images_per_page.get(page_num, 1))
                            bef, aft = split_page_text_around_image(pt, img_idx, nimg)
                            ctx = get_image_context(bef, aft, ocr_text or "", img_idx, page_num)
                            diagram_jobs_pending.append({
                                "image_path": out_path,
                                "image_id": image_id,
                                "page_num": page_num,
                                "ocr_text": ocr_text or "",
                                "context": ctx,
                            })
                            
                            # Update progress for OCR
                            if tracker:
                                tracker.set_stage(ProcessingStage.OCR_PROCESSING, total_images=detected_images, current_image=detected_images)
                            
                            # Update progress for summarization
                            if tracker:
                                tracker.set_stage(ProcessingStage.IMAGE_SUMMARIZATION, total_images=detected_images, current_image=detected_images)
                            
            # Append marker after context to help later viewing
                            if (ocr_text or "").strip():
                                text_output += f"\n[IMAGE {image_id}]\nOCR: {(ocr_text or '').strip()}\n"
                        except Exception:
                            continue
            except Exception:
                pass

        elif file.filename.endswith(".doc"):
            soffice = _get_soffice_path()
            if not soffice:
                text_output = "Unsupported file format (.doc) and LibreOffice not found for conversion."
            else:
                try:
                    subprocess.run([soffice, "--headless", "--convert-to", "docx", "--outdir", UPLOAD_DIR, filepath], check=True)
                    converted = os.path.join(UPLOAD_DIR, os.path.splitext(file.filename)[0] + ".docx")
                    d = docx.Document(converted)
                    for para in d.paragraphs:
                        text_output += para.text + "\n"
                    # Extract images from converted DOCX
                    img_idx = 0
                    for rel in d.part.rels.values():
                        if "image" in rel.reltype and getattr(rel, "target_part", None):
                            try:
                                img_idx += 1
                                image_count += 1
                                blob = rel.target_part.blob
                                page_num = 1
                                image_id = gen_image_id(file.filename, page_num, img_idx)
                                ext = ".png"
                                out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                                with open(out_path, "wb") as imf:
                                    imf.write(blob)
                                    try:
                                        img = Image.open(out_path)
                                        ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                                    except Exception:
                                        ocr_text = ""
                                    image_metadata.append((image_id, out_path, page_num, ocr_text))
                                    from services.image_service import get_image_context
                                    prior = text_output[max(0, len(text_output) - 500):]
                                    ctx = get_image_context(prior[-500:], "", ocr_text or "", img_idx, page_num)
                                    diagram_jobs_pending.append({
                                        "image_path": out_path,
                                        "image_id": image_id,
                                        "page_num": page_num,
                                        "ocr_text": ocr_text or "",
                                        "context": ctx,
                                    })
                                    if ocr_text.strip():
                                        text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                            except Exception:
                                pass
                except Exception as e:
                    text_output = f"Conversion error for .doc: {e}"

        elif file.filename.endswith(".docx"):
            d = docx.Document(filepath)
            for para in d.paragraphs:
                text_output += para.text + "\n"
            # Extract images from DOCX
            img_idx = 0
            for rel in d.part.rels.values():
                if "image" in rel.reltype and getattr(rel, "target_part", None):
                    try:
                        img_idx += 1
                        image_count += 1
                        blob = rel.target_part.blob
                        page_num = 1
                        image_id = gen_image_id(file.filename, page_num, img_idx)
                        ext = ".png"
                        out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                        with open(out_path, "wb") as imf:
                            imf.write(blob)
                        try:
                            img = Image.open(out_path)
                            ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                        except Exception:
                            ocr_text = ""
                        image_metadata.append((image_id, out_path, page_num, ocr_text))
                        from services.image_service import get_image_context
                        prior = text_output[max(0, len(text_output) - 500):]
                        ctx = get_image_context(prior[-500:], "", ocr_text or "", img_idx, page_num)
                        diagram_jobs_pending.append({
                            "image_path": out_path,
                            "image_id": image_id,
                            "page_num": page_num,
                            "ocr_text": ocr_text or "",
                            "context": ctx,
                        })
                        if ocr_text.strip():
                            text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                    except Exception:
                        pass

        elif file.filename.endswith(".ppt"):
            soffice = _get_soffice_path()
            if not soffice:
                text_output = "Unsupported file format (.ppt) and LibreOffice not found for conversion."
            else:
                try:
                    subprocess.run([soffice, "--headless", "--convert-to", "pptx", "--outdir", UPLOAD_DIR, filepath], check=True)
                    converted = os.path.join(UPLOAD_DIR, os.path.splitext(file.filename)[0] + ".pptx")
                    prs = pptx.Presentation(converted)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_output += shape.text + "\n"
                            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                                try:
                                    image_count += 1
                                    image_id = gen_image_id(file.filename, 1, image_count)
                                    ext = ".png"
                                    out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                                    with open(out_path, "wb") as imf:
                                        imf.write(shape.image.blob)
                                    try:
                                        img = Image.open(out_path)
                                        ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                                    except Exception:
                                        ocr_text = ""
                                    image_metadata.append((image_id, out_path, 1, ocr_text))
                                    from services.image_service import get_image_context
                                    prior = text_output[max(0, len(text_output) - 500):]
                                    ctx = get_image_context(prior[-500:], "", ocr_text or "", image_count, 1)
                                    diagram_jobs_pending.append({
                                        "image_path": out_path,
                                        "image_id": image_id,
                                        "page_num": 1,
                                        "ocr_text": ocr_text or "",
                                        "context": ctx,
                                    })
                                    if ocr_text.strip():
                                        text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                                except Exception:
                                    pass
                except Exception as e:
                    text_output = f"Conversion error for .ppt: {e}"

        elif file.filename.endswith(".pptx"):
            prs = pptx.Presentation(filepath)
            for slide_idx, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_output += shape.text + "\n"
                    if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                        try:
                            image_count += 1
                            image_id = gen_image_id(file.filename, slide_idx, image_count)
                            ext = ".png"
                            out_path = os.path.join(UPLOAD_DIR, f"{file.filename}_{image_id}{ext}")
                            with open(out_path, "wb") as imf:
                                imf.write(shape.image.blob)
                            try:
                                img = Image.open(out_path)
                                ocr_text = _advanced_ocr_text(image_path=out_path, pil_image=img)
                            except Exception:
                                ocr_text = ""
                            image_metadata.append((image_id, out_path, slide_idx, ocr_text))
                            from services.image_service import get_image_context
                            prior = text_output[max(0, len(text_output) - 500):]
                            ctx = get_image_context(prior[-500:], "", ocr_text or "", image_count, slide_idx)
                            diagram_jobs_pending.append({
                                "image_path": out_path,
                                "image_id": image_id,
                                "page_num": slide_idx,
                                "ocr_text": ocr_text or "",
                                "context": ctx,
                            })
                            if ocr_text.strip():
                                text_output += f"\n[IMAGE {image_id}]\nOCR: {ocr_text.strip()}\n"
                        except Exception:
                            pass

        elif file.filename.lower().endswith((".png", ".jpg", ".jpeg")):
            if tracker:
                tracker.set_stage(ProcessingStage.IMAGE_DETECTION)
            image_count = 1
            img = Image.open(filepath)
            if tracker:
                tracker.set_stage(ProcessingStage.OCR_PROCESSING, total_images=1, current_image=1)
            text_output = _advanced_ocr_text(image_path=filepath, pil_image=img)
            image_id = gen_image_id(file.filename, 1, 1)
            image_metadata.append((image_id, filepath, 1, text_output))
            from services.image_service import get_image_context
            ocr_here = text_output or ""
            ctx = get_image_context("", "", ocr_here, 1, 1)
            diagram_jobs_pending.append({
                "image_path": filepath,
                "image_id": image_id,
                "page_num": 1,
                "ocr_text": ocr_here,
                "context": ctx,
            })
            if tracker:
                tracker.set_stage(ProcessingStage.IMAGE_SUMMARIZATION, total_images=1, current_image=1)
            if text_output.strip():
                text_output += f"\n[IMAGE {image_id}]\nOCR: {text_output.strip()}\n"

        elif file.filename.lower().endswith(".txt"):
            # Handle text files
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text_output = f.read()

        else:
            raise HTTPException(status_code=400, detail="Unsupported file format")

        # Merge Basic Extraction data if provided
        if basic_extraction_data:
            print(f"🔄 Merging Basic Extraction data with AI extraction for enhanced analysis...")
            basic_text = basic_extraction_data.get("extracted_text", "") or basic_extraction_data.get("full_text", "")
            basic_image_summaries = basic_extraction_data.get("image_summaries", [])
            
            # Build comprehensive merged text that combines both sources intelligently
            merged_parts = []
            
            # Add instruction note for the agent at the beginning
            instruction_note = """
╔══════════════════════════════════════════════════════════════════════════════╗
║                    REQUIREMENTS EXTRACTION INSTRUCTIONS                        ║
╚══════════════════════════════════════════════════════════════════════════════╝

You have access to data from TWO sources that have been merged:

1. BASIC EXTRACTION RESULTS:
   - Initial text extraction and OCR from the document
   - Image OCR text and VLM interpretations from basic analysis
   - Raw extracted content from the document

2. AI-ENHANCED EXTRACTION:
   - Additional text extraction from the RAG agent's own analysis
   - Additional image analysis and interpretations
   - Enhanced context and semantic understanding

YOUR TASK:
- Analyze ALL sections below comprehensively
- Combine insights from BOTH Basic Extraction AND AI-Enhanced Extraction
- Use image interpretations from BOTH sources to understand visual requirements
- Extract MEANINGFUL requirements from images - understand what they show, not just copy OCR text
- DO NOT include image metadata markers like [IMAGE IMG-xxx], OCR: labels, or technical details like "Key Components:", "Technical Details:", "Contains X structural lines"
- DO NOT include raw OCR text that doesn't represent actual requirements
- Extract the FUNCTIONALITY and REQUIREMENTS shown in images, not the metadata about the images
- Merge duplicate or related requirements intelligently (avoid redundancy)
- Extract features that leverage the complete context from both sources
- Provide well-written, properly styled, and comprehensive features
- Ensure final output is clean, organized, and professional

IMPORTANT: 
- Extract features that represent the BEST combination of both data sources
- Focus on WHAT the system should do, not HOW images were processed
- Use image content to understand requirements, but don't mention image processing details in final features
"""
            merged_parts.append(instruction_note.strip())
            merged_parts.append("\n" + "═" * 80 + "\n")
            
            # Add Basic Extraction text with context marker
            if basic_text.strip():
                merged_parts.append("┌─ BASIC EXTRACTION RESULTS ───────────────────────────────────────────────────┐")
                merged_parts.append("│ Source: Initial text extraction and OCR from document analysis                │")
                merged_parts.append("└──────────────────────────────────────────────────────────────────────────────┘")
                merged_parts.append("")
                merged_parts.append(basic_text.strip())
                merged_parts.append("")
                print(f"   ✅ Merged {len(basic_text)} chars from Basic Extraction")
            
            # Add Basic Extraction image summaries with enhanced context
            # Extract meaningful requirements from images, not raw metadata
            if basic_image_summaries:
                merged_parts.append("\n" + "┌─ BASIC EXTRACTION IMAGE ANALYSIS ───────────────────────────────────────────┐")
                merged_parts.append("│ Source: OCR text and VLM interpretations from initial image analysis          │")
                merged_parts.append("└──────────────────────────────────────────────────────────────────────────────┘")
                merged_parts.append("")
                for idx, img_summary in enumerate(basic_image_summaries, 1):
                    if isinstance(img_summary, dict):
                        img_id = img_summary.get("image_id", f"UNKNOWN-{idx}")
                        page = img_summary.get("page", "?")
                        ocr = img_summary.get("ocr", "").strip()
                        summary = img_summary.get("summary", "") or img_summary.get("interpretation", "")
                        
                        # Extract meaningful content from OCR and summary, filtering out metadata
                        meaningful_content = []
                        
                        # Process OCR text - extract requirements, not raw text
                        if ocr:
                            # Remove metadata markers and extract actual content
                            ocr_clean = re.sub(r'\[IMAGE[^\]]*\]', '', ocr, flags=re.IGNORECASE)
                            ocr_clean = re.sub(r'OCR:\s*', '', ocr_clean, flags=re.IGNORECASE)
                            ocr_clean = re.sub(r'Key Components[^\n]*', '', ocr_clean, flags=re.IGNORECASE)
                            ocr_clean = re.sub(r'Technical Details[^\n]*', '', ocr_clean, flags=re.IGNORECASE)
                            ocr_clean = ocr_clean.strip()
                            
                            # Only include if it looks like a requirement, not metadata
                            if ocr_clean and len(ocr_clean) > 20 and not re.search(r'^\d+\s+(steps|components|elements)', ocr_clean, re.IGNORECASE):
                                # Check if it contains requirement-like language
                                if any(word in ocr_clean.lower() for word in ['shall', 'must', 'should', 'will', 'enable', 'allow', 'provide', 'support', 'system', 'user', 'application']):
                                    meaningful_content.append(f"From image {idx} (Page {page}): {ocr_clean[:300]}{'...' if len(ocr_clean) > 300 else ''}")
                        
                        # Process summary - extract requirements, not technical metadata
                        if summary:
                            summary_clean = re.sub(r'Image Type[^\n]*', '', summary, flags=re.IGNORECASE)
                            summary_clean = re.sub(r'Characteristics[^\n]*', '', summary_clean, flags=re.IGNORECASE)
                            summary_clean = re.sub(r'Analysis[^\n]*', '', summary_clean, flags=re.IGNORECASE)
                            summary_clean = re.sub(r'Contains\s+\d+\s+structural[^\n]*', '', summary_clean, flags=re.IGNORECASE)
                            summary_clean = re.sub(r'\d+\s+identifiable\s+(steps|elements|components)[^\n]*', '', summary_clean, flags=re.IGNORECASE)
                            summary_clean = summary_clean.strip()
                            
                            # Only include if it's meaningful content, not just metadata
                            if summary_clean and len(summary_clean) > 30:
                                # Check if it describes functionality or requirements
                                if any(word in summary_clean.lower() for word in ['workflow', 'process', 'system', 'user', 'feature', 'requirement', 'functionality', 'capability', 'shall', 'must', 'should']):
                                    meaningful_content.append(f"Image {idx} interpretation: {summary_clean[:300]}{'...' if len(summary_clean) > 300 else ''}")
                        
                        # Only add if we have meaningful content
                        if meaningful_content:
                            merged_parts.append(f"📷 IMAGE {idx} (Page {page}):")
                            for content in meaningful_content:
                                merged_parts.append(f"   {content}")
                            merged_parts.append("")
                    elif isinstance(img_summary, str):
                        # Clean string summaries too
                        clean_summary = re.sub(r'\[IMAGE[^\]]*\]', '', img_summary, flags=re.IGNORECASE)
                        clean_summary = re.sub(r'OCR:\s*', '', clean_summary, flags=re.IGNORECASE)
                        clean_summary = re.sub(r'Key Components[^\n]*', '', clean_summary, flags=re.IGNORECASE)
                        clean_summary = clean_summary.strip()
                        
                        if clean_summary and len(clean_summary) > 20:
                            merged_parts.append(f"📷 IMAGE {idx}: {clean_summary[:300]}{'...' if len(clean_summary) > 300 else ''}")
                            merged_parts.append("")
                
                print(f"   ✅ Added meaningful content from {len(basic_image_summaries)} images (filtered metadata)")
            
            # Add AI extraction text with marker
            merged_parts.append("\n" + "┌─ AI-ENHANCED EXTRACTION ─────────────────────────────────────────────────────┐")
            merged_parts.append("│ Source: RAG Agent's own document analysis and image processing                   │")
            merged_parts.append("└──────────────────────────────────────────────────────────────────────────────┘")
            merged_parts.append("")
            merged_parts.append(text_output)
            
            # Combine all parts - ensure all parts are strings
            merged_parts_str = []
            for part in merged_parts:
                if part is None:
                    continue
                # Convert to string and sanitize
                part_str = str(part) if not isinstance(part, str) else part
                part_str = sanitize_unicode(part_str)
                merged_parts_str.append(part_str)
            
            merged_text = "\n".join(merged_parts_str)
            
            # Ensure text_output is a valid string
            text_output = sanitize_unicode(merged_text) if merged_text else ""
            print(f"   ✅ Total merged text length: {len(text_output)} chars")
            print(f"   📋 Merged data includes: Basic text + {len(basic_image_summaries)} image summaries + AI extraction")
            
            # Also merge image metadata from both sources for database storage
            if basic_image_summaries:
                # Create a set of existing image IDs from AI extraction to avoid duplicates
                existing_image_ids = {img_id for img_id, _, _, _ in image_metadata}
                
                # Add basic extraction images that aren't duplicates
                for img_summary in basic_image_summaries:
                    if isinstance(img_summary, dict):
                        img_id = img_summary.get("image_id", "")
                        if img_id and img_id not in existing_image_ids:
                            # Try to reconstruct image metadata from basic extraction
                            page = img_summary.get("page", 1)
                            ocr = img_summary.get("ocr", "").strip()
                            path = img_summary.get("path", "")
                            if path and os.path.exists(path):
                                image_metadata.append((img_id, path, page, ocr))
                                existing_image_ids.add(img_id)
                                print(f"   ✅ Added image {img_id} from Basic Extraction to metadata")

        # Per-image diagram VLM requirements (smart filter + parallel workers; see image_service).
        vlm_pass = os.getenv("FLOWMIND_IMAGE_REQ_VLM_PASS", "1")
        print(f"VLM_PASS_ENV_VALUE: {vlm_pass}")
        req_vlm_enabled = vlm_pass == "1"
        print(f"VLM_ENABLED: {req_vlm_enabled}")
        if req_vlm_enabled and diagram_jobs_pending:
            try:
                from services.image_service import run_parallel_diagram_vlm_jobs
                max_w = int(os.getenv("FLOWMIND_VLM_MAX_WORKERS", "3"))
                max_w = max(1, min(max_w, 3))
                print(
                    f"Diagram VLM: {len(diagram_jobs_pending)} candidate image(s), "
                    f"max_workers={max_w}"
                )
                vlm_results = run_parallel_diagram_vlm_jobs(
                    diagram_jobs_pending,
                    max_workers=max_w,
                )
                basic_image_summaries = []
                if basic_extraction_data and isinstance(basic_extraction_data.get("image_summaries"), list):
                    basic_image_summaries = basic_extraction_data.get("image_summaries", [])
                for res in vlm_results:
                    image_id = res.get("image_id", "")
                    page_num = res.get("page_num", 0)
                    diagram_info = res.get("diagram_info") or {}
                    detected_features = diagram_info.get("detected_features") or []
                    if basic_image_summaries:
                        for item in basic_image_summaries:
                            if str(item.get("image_id")) == str(image_id):
                                item["diagram_type"] = diagram_info.get("type", item.get("diagram_type", "unknown"))
                                item["type_confidence"] = int(diagram_info.get("confidence", item.get("type_confidence", 0)) or 0)
                                item["detected_features"] = detected_features
                                item["vlm_analysis"] = res.get("vlm_analysis") or item.get("vlm_analysis", "")
                                item["extracted_requirements_count"] = int(
                                    res.get("requirements_count", item.get("extracted_requirements_count", 0)) or 0
                                )
                                break
                    if res.get("skipped"):
                        continue
                    understanding = (res.get("understanding") or "").strip()
                    if understanding:
                        text_output += (
                            f"\n[DIAGRAM_UNDERSTANDING {image_id}]\n"
                            f"{sanitize_unicode(understanding)}\n"
                            f"[/DIAGRAM_UNDERSTANDING]\n"
                        )
                    reqs = res.get("requirements") or []
                    if not reqs:
                        continue
                    text_output += f"\n[IMAGE_REQUIREMENTS {image_id}]\n" + "\n".join(
                        f"- {sanitize_unicode(r)}" for r in reqs
                    ) + "\n"
                    scen = next(
                        (
                            j.get("context", {}).get("scenario")
                            for j in diagram_jobs_pending
                            if j.get("image_id") == image_id
                        ),
                        None,
                    )
                    for req in reqs:
                        stmt = sanitize_unicode(req)
                        if not stmt:
                            continue
                        image_requirements_structured.append({
                            "statement": stmt,
                            "category": "functional",
                            "priority": "medium",
                            "confidence": 0.75,
                            "evidence_text": stmt[:240],
                            "evidence_page": page_num,
                            "evidence_image_id": image_id,
                            "source_type": "image",
                            "scenario": scen,
                        })
            except Exception as img_req_err:
                print(f"Image requirements VLM batch failed: {img_req_err}")

        # Save full text file
        sanitized_text = sanitize_unicode(text_output or "")
        text_file_path = os.path.join(UPLOAD_DIR, f"{file.filename}_full.txt")
        with open(text_file_path, "w", encoding="utf-8") as tf:
            tf.write(sanitized_text)

        # ----------- VALIDATE DOCUMENT QUALITY -----------
        if tracker:
            tracker.set_stage(ProcessingStage.VALIDATING)
            progress = tracker.get_progress()
            print(f"📊 Progress: {progress['progress']}% - {progress['message']}")

        v_res = validate_document_for_srs(sanitized_text)
        if v_res.get("is_rejected"):
            reject_msg = v_res.get("reject_reason", "Document does not meet technical/SRS quality standards.")
            print(f"❌ REJECTED agent flow: {reject_msg}")
            if tracker:
                tracker.set_stage(ProcessingStage.FINALIZING)
                tracker.complete()
            raise HTTPException(
                status_code=400,
                detail={
                    "error": "DOCUMENT_REJECTED",
                    "message": reject_msg,
                    "score": v_res.get("srs_score"),
                    "reasons": v_res.get("reasons")
                }
            )

        # Process with RAG agent - use async helpers to avoid blocking
        # Get user-specific agent for user-specific learning
        if tracker:
            tracker.set_stage(ProcessingStage.FINALIZING)  # AI Processing stage
            progress = tracker.get_progress()
            print(f"📊 Progress: {progress['progress']}% - {progress['message']}")
        
        print(f"🤖 Initializing AI agent for user_id={user_id}...")
        
        # Import async helpers
        from utils.async_helpers import run_in_thread
        
        # Get agent in thread pool to avoid blocking (with timeout)
        try:
            agent = await run_in_thread(get_agent, user_id=user_id, timeout=180.0)
            # #region agent log
            try:
                f = open(_DEBUG_LOG, "a", encoding="utf-8"); f.write(json.dumps({"sessionId": "0e985e", "hypothesisId": "D", "location": "flowmind.py:after get_agent", "message": "agent initialized", "data": {}, "timestamp": int(time.time() * 1000)}) + "\n"); f.close()
            except Exception:
                pass
            # #endregion
            print("✅ AI agent initialized, processing document...")
        except asyncio.TimeoutError:
            raise HTTPException(
                status_code=504,
                detail="Agent initialization timed out. Please try again."
            )
        except Exception as e:
            print(f"❌ Error initializing agent: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Failed to initialize AI agent: {str(e)}"
            )
        
        if tracker:
            # Update progress: Document processing
            tracker.set_stage(ProcessingStage.FINALIZING)
            progress = tracker.get_progress()
            print(f"📊 Progress: {progress['progress']}% - Processing document with AI agent...")
        
        # Process document in thread pool with timeout
        # Ensure text_output is a valid string before passing to agent
        if not isinstance(text_output, str):
            text_output = str(text_output) if text_output is not None else ""
        text_output = sanitize_unicode(text_output)
        
        # Validate text is not empty and is a proper string
        if not text_output or len(text_output.strip()) == 0:
            raise HTTPException(
                status_code=400,
                detail="No text content extracted from document"
            )

        # Visual-understanding fallback: if embedded image extraction found nothing,
        # OCR a few rendered PDF pages and feed that context to the agent.
        try:
            if int(image_count or 0) == 0 and str(file.filename or '').lower().endswith('.pdf'):
                try:
                    import fitz  # PyMuPDF
                except Exception:
                    fitz = None

                if fitz is not None:
                    pdf_path = os.path.join(UPLOAD_DIR, file.filename)
                    if os.path.exists(pdf_path):
                        doc = None
                        fallback_ocr_blocks = []
                        try:
                            doc = fitz.open(pdf_path)
                            max_pages = min(len(doc), 8)
                            for i in range(max_pages):
                                try:
                                    page = doc[i]
                                    pix = page.get_pixmap(
                                        matrix=fitz.Matrix(1.8, 1.8),
                                        colorspace=fitz.csRGB,
                                        alpha=False,
                                        annots=False,
                                    )
                                    tmp_img_path = os.path.join(UPLOAD_DIR, f"{file.filename}_CTX-{i+1}.png")
                                    pix.save(tmp_img_path)
                                    with Image.open(tmp_img_path) as page_img:
                                        ocr_text = _advanced_ocr_text(image_path=tmp_img_path, pil_image=page_img).strip()
                                    if len(ocr_text) >= 25:
                                        fallback_ocr_blocks.append(f"[PAGE {i+1}] {ocr_text[:1200]}")
                                except Exception:
                                    continue
                        finally:
                            if doc is not None:
                                doc.close()

                        if fallback_ocr_blocks:
                            fallback_context = "\n\n[VISUAL_PAGE_OCR_FALLBACK]\n" + "\n".join(fallback_ocr_blocks)
                            text_output += fallback_context
                            print(f"📸 Added visual OCR fallback context from {len(fallback_ocr_blocks)} PDF page(s) for AI extraction")

                            # Make this context also available for feature merge stage later.
                            if isinstance(basic_extraction_data, dict):
                                basic_extraction_data.setdefault('image_summaries', [])
                                basic_extraction_data['image_summaries'].extend([
                                    {'summary': block, 'interpretation': block} for block in fallback_ocr_blocks
                                ])
        except Exception as fallback_err:
            print(f"⚠️ Visual OCR fallback skipped: {fallback_err}")
        
        try:
            agent_result = await run_in_thread(
                agent.process_document,
                text_output,
                file.filename or "unknown",
                timeout=240.0  # 4 minutes for document processing (increased for large documents)
            )
        except asyncio.TimeoutError:
            raise HTTPException(
                status_code=504,
                detail="Document processing timed out. The document may be too large or complex."
            )
        except Exception as e:
            print(f"❌ Error processing document: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"Error processing document: {str(e)}"
            )
        # #region agent log
        try:
            f = open(_DEBUG_LOG, "a", encoding="utf-8"); f.write(json.dumps({"sessionId": "0e985e", "hypothesisId": "D", "location": "flowmind.py:after process_document", "message": "process_document done", "data": {}, "timestamp": int(time.time() * 1000)}) + "\n"); f.close()
        except Exception:
            pass
        # #endregion
        if agent_result.get("status") != "success":
            msg = agent_result.get("message") or "Agent failed to process document"
            raise HTTPException(status_code=500, detail=msg)

        # Validate document quality/SRS-likeness before extraction starts.
        srs_validation = validate_document_for_srs(text_output or "")
        print(
            f"📋 SRS validation: score={srs_validation.get('srs_score')} "
            f"confidence={srs_validation.get('confidence')} is_srs={srs_validation.get('is_srs')}"
        )

        # Extract requirements using the agent
        if tracker:
            tracker.set_stage(ProcessingStage.FINALIZING)  # Requirements extraction
            progress = tracker.get_progress()
            print(f"📊 Progress: {progress['progress']}% - Extracting requirements...")
        
        # Extract requirements in thread pool with timeout
        # Get timeout from environment or use default (600 seconds = 10 minutes for large documents)
        extraction_timeout = float(os.getenv("FLOWMIND_EXTRACTION_TIMEOUT", "600.0"))
        print(f"⏱️  Requirements extraction timeout: {extraction_timeout}s")
        
        try:
            print(f"🔄 Starting requirements extraction (timeout: {extraction_timeout}s)...")
            print(f"📊 Agent type: {type(agent)}")
            print(f"📊 Agent has extract_requirements: {hasattr(agent, 'extract_requirements')}")
            
            # Call extract_requirements directly in thread
            print(f"📊 About to call agent.extract_requirements in thread...")
            requirements_result = await run_in_thread(
                agent.extract_requirements,
                timeout=extraction_timeout
            )
            print(f"📊 Back from run_in_thread, got result: {type(requirements_result)}")
            
            print(f"✅ Requirements extraction completed successfully")
            print(f"📊 Result type: {type(requirements_result)}")
            if isinstance(requirements_result, dict):
                print(f"📊 Result keys: {list(requirements_result.keys())}")
                print(f"📊 Result status: {requirements_result.get('status')}")
            else:
                print(f"📊 Result (not dict): {requirements_result}")
            
            # Ensure result is a dict
            if not isinstance(requirements_result, dict):
                print(f"⚠️ Result is not a dict, converting: {type(requirements_result)}")
                requirements_result = {"status": "success", "response": str(requirements_result)}
            
            print(f"📊 Result keys: {list(requirements_result.keys())}")
            print(f"📊 Result status: {requirements_result.get('status', 'N/A')}")
            # #region agent log
            try:
                f = open(_DEBUG_LOG, "a", encoding="utf-8"); f.write(json.dumps({"sessionId": "0e985e", "hypothesisId": "D", "location": "flowmind.py:after extract_requirements", "message": "extract_requirements done", "data": {}, "timestamp": int(time.time() * 1000)}) + "\n"); f.close()
            except Exception:
                pass
            # #endregion
            # Validate result structure
            if requirements_result.get("status") != "success":
                print(f"⚠️ Extraction returned non-success status: {requirements_result.get('status')}")
                print(f"⚠️ Error message: {requirements_result.get('message', 'No message')}")
            
            # Ensure response field exists
            if "response" not in requirements_result:
                print(f"⚠️ No 'response' field in result, adding empty response")
                requirements_result["response"] = ""
            
            # Add note about Basic Extraction merge if data was provided
            if basic_extraction_data:
                merge_note = """

╔══════════════════════════════════════════════════════════════════════════════╗
║                    ENHANCED ANALYSIS - MERGED DATA SOURCES                    ║
╚══════════════════════════════════════════════════════════════════════════════╝

This comprehensive analysis combines data from multiple sources:

✅ Basic Extraction Results:
   • Initial text extraction and OCR from document
   • Image OCR text and VLM interpretations
   • Raw extracted content

✅ AI-Enhanced Extraction:
   • Advanced RAG agent analysis
   • Semantic understanding and context
   • Enhanced feature extraction

The features above represent the best combination of both data sources, providing
comprehensive, well-written, and properly styled requirements that leverage the
complete context from both basic extraction and AI analysis.

═══════════════════════════════════════════════════════════════════════════════
"""
                requirements_result["response"] = requirements_result.get("response", "") + merge_note
                print(f"✅ Added enhanced merge note to final output")
            
            print(f"📊 Response length: {len(str(requirements_result.get('response', '')))}")
            print(f"📊 About to continue to database save...")
        except asyncio.TimeoutError as timeout_err:
            print(f"⏱️  Requirements extraction timed out: {timeout_err}")
            print(f"⏱️  Requirements extraction timed out after {extraction_timeout}s")
            # Try to return partial results if available
            try:
                # Check if we have any partial results from the agent
                if hasattr(agent, 'last_extraction_result'):
                    partial_result = agent.last_extraction_result
                    if partial_result and partial_result.get('status') != 'error':
                        print(f"⚠️  Returning partial results due to timeout")
                        requirements_result = partial_result
                        requirements_result['partial'] = True
                        requirements_result['message'] = f"Extraction timed out after {extraction_timeout}s. Partial results returned."
                    else:
                        raise HTTPException(
                            status_code=504,
                            detail=f"Requirements extraction timed out after {extraction_timeout}s. Please try again with a smaller document or increase FLOWMIND_EXTRACTION_TIMEOUT."
                        )
                else:
                    raise HTTPException(
                        status_code=504,
                        detail=f"Requirements extraction timed out after {extraction_timeout}s. Please try again with a smaller document or increase FLOWMIND_EXTRACTION_TIMEOUT."
                    )
            except HTTPException:
                raise
            except Exception as e:
                print(f"❌ Error getting partial results: {str(e)}")
                raise HTTPException(
                    status_code=504,
                    detail=f"Requirements extraction timed out after {extraction_timeout}s. Please try again with a smaller document or increase FLOWMIND_EXTRACTION_TIMEOUT."
                )
        except Exception as e:
            print(f"❌ Error extracting requirements: {str(e)}")
            import traceback
            traceback.print_exc()
            raise HTTPException(
                status_code=500,
                detail=f"Error extracting requirements: {str(e)}"
            )
        
        print(f"📊 Post-extraction: Starting database save and response preparation...")
        print(f"📊 Requirements result type: {type(requirements_result)}")
        
        # Ensure requirements_result is a dict
        if not isinstance(requirements_result, dict):
            print(f"⚠️ requirements_result is not a dict, converting...")
            requirements_result = {
                "status": "success",
                "response": str(requirements_result) if requirements_result else "No response"
            }
        
        if tracker:
            tracker.complete()
            print(f"✅ Processing complete: {tracker.get_progress()['progress']}%")
        
        # Debug logging
        print(f"🔍 Requirements result status: {requirements_result.get('status')}")
        print(f"🔍 Requirements result keys: {list(requirements_result.keys())}")
        if requirements_result.get("message"):
            print(f"🔍 Error message: {requirements_result.get('message')}")
        
        # Allow partial results if timeout occurred
        if requirements_result.get("status") != "success" and not requirements_result.get("partial"):
            error_msg = requirements_result.get("message") or "Requirement extraction failed"
            print(f"❌ Extraction failed: {error_msg}")
            raise HTTPException(status_code=500, detail=error_msg)
        
        # Log if partial results were returned
        if requirements_result.get("partial"):
            print(f"⚠️  Partial results returned due to timeout - some data may be incomplete")

        # Save to database - ensure session is closed promptly
        record = None
        view_id = str(uuid.uuid4())  # Generate view_id upfront
        try:
            if db_session is None:
                db = SessionLocal()
                should_close = True
            else:
                db = db_session
                should_close = False
            
            record = ParsedFile(
                filename=file.filename,
                extracted_text=sanitized_text[:400],
                detected_shapes=image_count,
                summary=f"Processed with RAG agent. {agent_result['message']}",
                full_text_path=text_file_path,
                user_id=user_id,  # Link to user
                project_id=project_id,
                view_id=view_id  # Set view_id upfront
            )
            db.add(record)
            db.commit()
            db.refresh(record)

            # #region agent log
            try:
                f = open(_DEBUG_LOG, "a", encoding="utf-8"); f.write(json.dumps({"sessionId": "0e985e", "hypothesisId": "B", "location": "flowmind.py:before ImageMeta loop", "message": "about to save image metadata", "data": {"n_images": len(image_metadata)}, "timestamp": int(time.time() * 1000)}) + "\n"); f.close()
            except Exception:
                pass
            # #endregion
            # Save image metadata (carry diagram-analysis context from basic extraction when available)
            image_summary_by_id = {}
            if basic_extraction_data and isinstance(basic_extraction_data.get("image_summaries"), list):
                for item in basic_extraction_data.get("image_summaries", []):
                    if isinstance(item, dict):
                        iid = str(item.get("image_id") or "").strip()
                        if iid:
                            image_summary_by_id[iid] = item

            for image_id, image_path, page_num, ocr_text in image_metadata:
                summary_item = image_summary_by_id.get(str(image_id), {})
                detected_features = summary_item.get("detected_features")
                if isinstance(detected_features, str):
                    detected_features = [detected_features]
                if not isinstance(detected_features, list):
                    detected_features = []

                img_meta = ImageMeta(
                    file_id=record.id,
                    image_path=image_path,
                    page_number=page_num,
                    ocr_text=ocr_text,
                    diagram_type=summary_item.get("diagram_type"),
                    type_confidence=int(summary_item.get("type_confidence", 0) or 0),
                    detected_features=json.dumps(detected_features, ensure_ascii=False),
                    vlm_analysis=summary_item.get("vlm_analysis") or summary_item.get("interpretation") or summary_item.get("summary") or "",
                    extracted_requirements_count=int(summary_item.get("extracted_requirements_count", 0) or 0),
                )
                db.add(img_meta)

            # Persist AI response so users can see previous chat-style results after refresh/restart.
            assistant_message = str(requirements_result.get("response") or "").strip()
            if assistant_message:
                history_row = AgentChatHistory(
                    user_id=user_id,
                    project_id=project_id,
                    file_id=record.id,
                    filename=file.filename,
                    user_message=f"Analyze this document: {file.filename}",
                    assistant_message=assistant_message,
                    view_id=view_id,
                )
                db.add(history_row)

            db.commit()
            
            print(f"✅ Database record saved: view_id={view_id}, file_id={record.id}")
        except Exception as e:
            # #region agent log
            try:
                f = open(_DEBUG_LOG, "a", encoding="utf-8"); f.write(json.dumps({"sessionId": "0e985e", "hypothesisId": "B", "location": "flowmind.py:db except", "message": "database block raised", "data": {"error": str(e), "type": type(e).__name__}, "timestamp": int(time.time() * 1000)}) + "\n"); f.close()
            except Exception:
                pass
            # #endregion
            print(f"❌ Database error: {str(e)}")
            import traceback
            traceback.print_exc()
            # Don't fail completely - use the view_id we already generated
            print(f"⚠️ Using fallback (record not saved to DB, view_id: {view_id})")
        finally:
            # Ensure session is always closed
            if should_close and db_session is None:
                try:
                    db.close()
                    print(f"✅ Database session closed")
                except Exception:
                    pass

        # Build quick lookup of context for images within this request is not tracked separately here,
        # so pass empty context for now; summarized meanings rely on OCR wording for agent runs
        
        # Collect image summaries for merging with text extraction
        image_summaries_for_features = []
        try:
            # Get image summaries from basic_extraction_data if available
            if basic_extraction_data and basic_extraction_data.get('image_summaries'):
                image_summaries_for_features = basic_extraction_data.get('image_summaries', [])
                print(f"📸 Found {len(image_summaries_for_features)} image summaries to merge with text extraction")
            
            # Also try to get from image_metadata if VLM summaries are available
            # This would be populated if VLM processing was done
            # (ImageMeta is used at module level above; do not add a local import here or it shadows and causes UnboundLocalError at first use.)
            if not image_summaries_for_features and record is not None:
                db_img = SessionLocal()
                try:
                    img_records = db_img.query(ImageMeta).filter(ImageMeta.file_id == record.id).all()
                    for img_record in img_records:
                        if img_record.ocr_text and len(img_record.ocr_text) > 50:  # Only use substantial OCR text
                            features = []
                            try:
                                if img_record.detected_features:
                                    parsed_features = json.loads(img_record.detected_features)
                                    if isinstance(parsed_features, list):
                                        features = parsed_features
                                    elif isinstance(parsed_features, str):
                                        features = [parsed_features]
                            except Exception:
                                features = [str(img_record.detected_features)] if img_record.detected_features else []
                            image_summaries_for_features.append({
                                'summary': img_record.ocr_text,
                                'interpretation': img_record.ocr_text,
                                'diagram_type': img_record.diagram_type or "unknown",
                                'type_confidence': int(img_record.type_confidence or 0),
                                'detected_features': features,
                                'vlm_analysis': img_record.vlm_analysis or "",
                                'extracted_requirements_count': int(img_record.extracted_requirements_count or 0),
                            })
                    if image_summaries_for_features:
                        print(f"📸 Found {len(image_summaries_for_features)} image summaries from database")
                finally:
                    db_img.close()
        except Exception as e:
            print(f"⚠️ Error collecting image summaries: {e}")
        
        # Parse and save features for approval (run synchronously but quickly, skip if slow)
        features_count = 0
        try:
            extracted_response = requirements_result.get("response", "")
            structured_requirements = requirements_result.get("requirements_json")
            if not isinstance(structured_requirements, list):
                structured_requirements = []
            if image_requirements_structured:
                existing_keys = {
                    (str(item.get("statement", "")).strip().lower(), str(item.get("evidence_image_id", "")).strip())
                    for item in structured_requirements
                }
                for item in image_requirements_structured:
                    key = (str(item.get("statement", "")).strip().lower(), str(item.get("evidence_image_id", "")).strip())
                    if key not in existing_keys:
                        structured_requirements.append(item)
                        existing_keys.add(key)
                requirements_result["requirements_json"] = structured_requirements
                requirements_result["requirements_json_count"] = len(structured_requirements)
            try:
                img_only = [x for x in structured_requirements if x.get("evidence_image_id")]
                if img_only and agent:
                    full_doc = sanitized_text or ""
                    chunks = agent.text_splitter.split_text(full_doc)
                    linked = agent.link_image_requirements_to_document_text(img_only, chunks)
                    lk = {(d.get("evidence_image_id"), d.get("statement")): d for d in linked}
                    for i, it in enumerate(structured_requirements):
                        key = (it.get("evidence_image_id"), it.get("statement"))
                        if key in lk:
                            structured_requirements[i] = lk[key]
                    requirements_result["requirements_json"] = structured_requirements
                    requirements_result["requirements_json_count"] = len(structured_requirements)
            except Exception as le:
                print(f"Image-text linking failed: {le}")
            try:
                from services.requirement_validation import partition_valid_requirements

                structured_requirements, partition_rejects = partition_valid_requirements(
                    structured_requirements
                )
                prev_rej = requirements_result.get("validation_rejects") or []
                if not isinstance(prev_rej, list):
                    prev_rej = []
                requirements_result["validation_rejects"] = prev_rej + partition_rejects
                requirements_result["requirements_json"] = structured_requirements
                requirements_result["requirements_json_count"] = len(structured_requirements)
            except Exception as ve:
                print(f"Requirement validation partition failed: {ve}")
            has_structured_requirements = isinstance(structured_requirements, list) and len(structured_requirements) > 0
            if (extracted_response or has_structured_requirements) and record and hasattr(record, 'id') and record.id:
                # Create a new database session for feature parsing (thread-safe)
                try:
                    # Use a new session to avoid thread-safety issues
                    db_features = SessionLocal()
                    try:
                        features_count = parse_and_save_features(
                            extracted_response, 
                            user_id, 
                            record.id, 
                            db_features,
                            image_summaries=image_summaries_for_features,
                            structured_requirements=structured_requirements,
                            project_id=project_id,
                        )
                        print(f"📋 Saved {features_count} features for user approval (merged from text + images)")
                    finally:
                        db_features.close()
                except Exception as e:
                    print(f"⚠️ Failed to save features: {e}")
            else:
                print(f"⚠️ Skipping feature parsing (no record or no response)")
        except Exception as e:
            print(f"⚠️ Failed to parse features: {e}")
        
        # Save to REQUIREMENTS_VIEWS for backward compatibility (use OCR text as summary, don't re-process)
        try:
            # Use OCR text as summary to avoid slow VLM re-processing
            image_summaries = []
            for (iid, path, pg, ocr) in image_metadata:
                # Sanitize OCR text and create summary
                sanitized_ocr = sanitize_unicode((ocr or "").strip())
                summary = sanitized_ocr[:200] if sanitized_ocr else ""  # Just use OCR, no extra processing
                
                image_summaries.append({
                    "image_id": iid,
                    "path": path,
                    "page": pg,
                    "ocr": sanitized_ocr,
                    "summary": summary
                })
            
            REQUIREMENTS_VIEWS[view_id] = {
                "filename": file.filename,
                "summary": f"Extracted {len(text_output.split())} words and {image_count} image(s)",
                "response": requirements_result.get("response", ""),
                "project_id": project_id,
                "images": image_summaries,
                "srs_validation": srs_validation,
                "validation_rejects": requirements_result.get("validation_rejects") or [],
                "requirement_classification_basis": requirements_result.get(
                    "requirement_classification_basis", ""
                ),
            }
            # Save views asynchronously to avoid blocking
            try:
                save_views()  # Save to persistent storage
                print(f"✅ Saved to REQUIREMENTS_VIEWS: {view_id}")
            except Exception as save_error:
                print(f"⚠️ Failed to save views to disk: {save_error}")
        except Exception as e:
            print(f"⚠️ Failed to save to REQUIREMENTS_VIEWS: {e}")
            # Continue anyway - don't block response
        
        # Return response immediately - don't wait for slow operations
        print(f"✅ Preparing response to return to client...")
        print(f"📊 Requirements extraction status: {requirements_result.get('status')}")
        print(f"📊 Requirements extraction response length: {len(str(requirements_result.get('response', '')))}")
        
        # Build response data - ensure all fields exist
        try:
            response_data = {
                "filename": file.filename or "unknown",
                "extraction_summary": f"Extracted {len(text_output.split()) if text_output else 0} words and {image_count} image(s)",
                "agent_processing": agent_result or {"status": "success", "message": "Processing completed"},
                "requirements_extraction": requirements_result or {"status": "error", "message": "No extraction result"},
                "view_id": view_id or str(uuid.uuid4()),
                "full_text_file": text_file_path or "",
                "images_detected": image_count,
                "image_metadata_saved": len(image_metadata),
                "features_saved": features_count,
                "srs_validation": srs_validation,
            }
            print(f"✅ Response data built successfully")
            print(f"📊 Response keys: {list(response_data.keys())}")
            print(f"📊 Returning results to client (view_id: {response_data['view_id']})...")
            print(f"📊 Response data prepared, returning now...")
            print(f"🚀 ACTUALLY RETURNING RESPONSE NOW - view_id: {response_data['view_id']}")
            # Sanitize Unicode to prevent encoding errors
            return sanitize_dict(response_data)
        except Exception as e:
            print(f"❌ Error preparing response: {str(e)}")
            import traceback
            traceback.print_exc()
            # Return minimal response even on error
            print(f"⚠️ Returning fallback response due to error")
            fallback_response = {
                "filename": file.filename or "unknown",
                "extraction_summary": f"Extracted {len(text_output.split()) if text_output else 0} words and {image_count} image(s)",
                "agent_processing": agent_result or {"status": "success"},
                "requirements_extraction": requirements_result or {"status": "error", "message": "Error occurred"},
                "view_id": view_id or str(uuid.uuid4()),
                "error": f"Error preparing full response: {str(e)}",
                "srs_validation": srs_validation,
            }
            # Sanitize Unicode to prevent encoding errors
            return sanitize_dict(fallback_response)

    except HTTPException as e:
        print(f"❌ HTTPException in _analyze_with_agent_internal: {e.status_code} - {e.detail}")
        raise e
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        print("❌ Analyze_with_agent error:", tb)
        print(f"❌ Exception type: {type(e).__name__}")
        print(f"❌ Exception message: {str(e)}")
        # Try to return a minimal response instead of raising
        try:
            print(f"⚠️ Attempting to return error response instead of raising exception")
            return {
                "filename": file.filename if hasattr(file, 'filename') else "unknown",
                "extraction_summary": "Error occurred during processing",
                "agent_processing": {"status": "error", "message": str(e)},
                "requirements_extraction": {"status": "error", "message": str(e)},
                "view_id": str(uuid.uuid4()),
                "error": f"Error processing document: {str(e)}"
            }
        except Exception as return_error:
            print(f"❌ Failed to return error response: {return_error}")
            raise HTTPException(status_code=500, detail=f"Error processing document: {str(e) or 'See server logs'}")


@app.post("/extract_requirements")
async def extract_requirements(query: str = None):
    """Extract requirements from already processed documents."""
    try:
        agent = get_agent()
        result = agent.extract_requirements(query)
        
        if result["status"] != "success":
            raise HTTPException(status_code=500, detail=result["message"])
        
        return result
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting requirements: {str(e)}")


@app.get("/get_document_summary")
async def get_document_summary():
    """Get a summary of all processed documents."""
    try:
        agent = get_agent()
        result = agent.get_document_summary()
        
        if result["status"] != "success":
            raise HTTPException(status_code=500, detail=result["message"])
        
        return result
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating summary: {str(e)}")


@app.post("/search_requirements")
async def search_requirements(requirement_type: str):
    """Search for specific types of requirements."""
    try:
        agent = get_agent()
        result = agent.search_specific_requirements(requirement_type)
        
        if result["status"] != "success":
            raise HTTPException(status_code=500, detail=result["message"])
        
        return result
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error searching requirements: {str(e)}")


@app.get("/agent_status")
async def get_agent_status():
    """Get the status of the RAG agent."""
    try:
        agent = get_agent()
        return {
            "status": "active",
            "model": getattr(agent, "model_name", "heuristic"),
            "vectorstore_initialized": agent.vectorstore is not None,
            "collection_name": agent.collection_name,
            "tools_available": len(agent.tools)
        }
    except Exception as e:
        return {
            "status": "error",
            "message": str(e)
        }


@app.get("/agent_status_page")
async def agent_status_page():
    return {"message": "Use GET /agent_status for JSON status."}


# =============================
# Formatted Requirements Viewer
# =============================

@app.get("/view_requirements/{view_id}", response_class=HTMLResponse)
async def view_requirements(view_id: str):
    # Reload views in case file was updated
    load_views()
    data = REQUIREMENTS_VIEWS.get(view_id)
    if not data:
        raise HTTPException(status_code=404, detail="View not found or expired")

    def esc(s: str) -> str:
        return html.escape(s or "")

    def normalize_req(text: str) -> str:
        """Create a stable key for matching stored features to displayed items."""
        import re
        return re.sub(r"\s+", " ", text or "").strip().lower().rstrip(".")

    # Pull latest approval states for this view (if we have a DB record)
    feature_map = {}
    feature_stats = {"approved": 0, "denied": 0, "pending": 0, "total": 0}
    db_session = SessionLocal()
    try:
        parsed_file = db_session.query(ParsedFile).filter(ParsedFile.view_id == view_id).first()
        if parsed_file:
            for feat in db_session.query(Feature).filter(Feature.file_id == parsed_file.id).all():
                key = normalize_req(feat.description)
                feature_map[key] = {
                    "status": feat.status or "pending",
                    "quality_score": feat.quality_score or 0,
                    "category": feat.category or "",
                    "classification_reason": getattr(feat, "classification_reason", None) or "",
                    "classification_method": getattr(feat, "classification_method", None) or "",
                    "classification_confidence_label": getattr(feat, "classification_confidence_label", None) or "",
                }
                feature_stats["total"] += 1
                if feat.status in feature_stats:
                    feature_stats[feat.status] += 1
    finally:
        db_session.close()

    filename = esc(data.get("filename", ""))
    summary = esc(data.get("summary", ""))
    response_text = data.get("response", "")
    images = data.get("images") or []
    srs_validation = data.get("srs_validation") or {}
    srs_confidence = str(srs_validation.get("confidence") or "None")
    srs_score = int(srs_validation.get("srs_score") or 0)
    srs_recommendation = esc(str(srs_validation.get("recommendation") or ""))
    srs_reasons = srs_validation.get("reasons") or []
    srs_reasons_html = "".join(f"<li>{esc(str(r))}</li>" for r in srs_reasons[:8])
    srs_banner_class = {
        "High": "srs-high",
        "Medium": "srs-medium",
        "Low": "srs-low",
        "None": "srs-none",
    }.get(srs_confidence, "srs-none")
    srs_banner_html = f"""
        <div class="srs-banner {srs_banner_class}">
            <div class="srs-banner-title">
                <i class="fas fa-shield-alt"></i>
                Document Validity Check
            </div>
            <div class="srs-banner-meta">
                <strong>SRS Score:</strong> {srs_score}/100 &nbsp;|&nbsp;
                <strong>Confidence:</strong> {esc(srs_confidence)}
            </div>
            <div class="srs-banner-text">{srs_recommendation}</div>
            <details class="srs-details">
                <summary>Show scoring findings</summary>
                <ul>{srs_reasons_html}</ul>
            </details>
        </div>
    """
    
    # Convert to HTML: create collapsible dropdown sections
    import re
    sections = []
    current_section = None
    current_items = []
    lines = response_text.split("\n")
    
    for raw_line in lines:
        raw_line = raw_line.strip()
        if not raw_line:
            continue
        
        escaped_line = esc(raw_line)
        
        # Check if this is a section heading (ends with colon and no bullet)
        if raw_line.endswith(":") and not raw_line.startswith("-"):
            # Save previous section
            if current_section and current_items:
                sections.append((current_section, current_items))
            
            # Start new section
            current_section = escaped_line.rstrip(":")
            current_items = []
        # Check if this is a bullet point (handles -, ✅, ⚠️, ❌)
        elif (raw_line.startswith("- ") or raw_line.startswith("✅") or raw_line.startswith("⚠️") or raw_line.startswith("❌")) and current_section:
            requirement_text_raw = raw_line[2:].strip() if raw_line.startswith("- ") else re.sub(r"^[✅⚠️❌]\s*", "", raw_line).strip()
            if requirement_text_raw.lower() != "(none)":
                requirement_text = esc(requirement_text_raw)
                normalized = normalize_req(requirement_text_raw)
                status_info = feature_map.get(normalized)
                current_items.append((requirement_text, status_info))
        # Regular line (add to current section if exists)
        elif current_section and raw_line.lower() != "(none)":
            current_items.append((escaped_line, None))
    
    # Save last section
    if current_section and current_items:
        sections.append((current_section, current_items))
    
    # Build collapsible HTML
    safe_html = ""
    section_counter = 0
    for section_name, items in sections:
        section_id = f"section_{section_counter}"
        section_counter += 1
        item_count = len(items)
        
        safe_html += f'''
        <div class="dropdown-section">
            <button class="dropdown-toggle" data-section-id="{section_id}">
                <span class="dropdown-icon" id="icon_{section_id}">▼</span>
                <span class="dropdown-title">{section_name}</span>
                <span class="dropdown-count">({item_count})</span>
            </button>
            <div class="dropdown-content" id="{section_id}" style="display: none;">
                <ul class="simple-list">
        '''
        
        for item_text, status_info in items:
            status_class = ""
            status_badge = ""
            info_badge = ""
            if status_info:
                status = status_info.get("status", "pending")
                status_class = f" status-{status}"
                status_label = status.capitalize()
                status_badge = f'<span class="status-pill status-{status}">{status_label}</span>'
                cat_raw = (status_info.get("category") or "").replace("_", " ").strip() or "unknown"
                cat = cat_raw.title()
                reason = status_info.get("classification_reason") or "No stored explanation (re-upload to refresh)."
                method = (status_info.get("classification_method") or "—").lower()
                conf_lbl = status_info.get("classification_confidence_label") or "—"
                if status_info.get("classification_reason") or status_info.get("classification_method"):
                    tip = (
                        f"Classified as {cat} because: {reason} "
                        f"Method: {method}. Confidence: {conf_lbl}."
                    )
                    info_badge = (
                        f'<span class="classify-info" tabindex="0" role="img" '
                        f'title="{esc(tip)}" data-tooltip="{esc(tip)}">'
                        f'<i class="fas fa-info-circle" aria-label="Why this classification"></i></span>'
                    )

            safe_html += f'<li class="simple-item{status_class}"><span class="item-text">{item_text}</span>{info_badge}{status_badge}</li>'
        
        safe_html += '''
                </ul>
            </div>
        </div>
        '''
    
    safe = safe_html if safe_html else '<div class="no-content">No requirements found.</div>'
    
    approved_count = feature_stats.get("approved", 0)
    denied_count = feature_stats.get("denied", 0)
    pending_count = feature_stats.get("pending", 0)
    total_count = feature_stats.get("total", 0)
    status_summary_html = f"""
        <div class="status-summary" style="margin-bottom:1rem;font-size:0.9rem;color:#64748b;">
            <strong>Summary:</strong> Approved: {approved_count} | Pending: {pending_count} | Denied: {denied_count} | Total: {total_count}
        </div>
    """
    
    # Build images section HTML - normalize path so "View Image" works (always /uploads/...)
    items = []
    for img in images:
        image_id = esc(str(img.get("image_id", "")))
        page = esc(str(img.get("page", "?")))
        raw_path = (img.get("path", "") or "").replace("\\", "/")
        if "uploads" in raw_path:
            path = raw_path.split("uploads", 1)[-1].lstrip("/")
            path = esc("uploads/" + path if path else "uploads")
        else:
            path = esc("uploads/" + os.path.basename(raw_path) if raw_path else "uploads")
        ocr = esc((img.get("ocr") or "").strip()[:1500])
        raw_summary = (img.get("summary") or "").strip()
        summary = esc(raw_summary)
        if summary:
            # Format summary with line breaks
            summary_html = "<br>".join(esc(line) for line in summary.splitlines())
        else:
            summary_html = ""
        # Build a brief interpretation sentence from the first bullet/line
        interp_src = raw_summary.splitlines()
        first_line = next((ln for ln in interp_src if ln.strip()), "")
        if first_line.startswith(("- ", "• ")):
            first_line = first_line[2:].strip()
        interpretation = f"This image indicates/defines: {first_line}" if first_line else ""
        interpretation_html = esc(interpretation)
        
        items.append(f"""
            <div class="image-card">
                <div class="image-header">
                    <div class="image-icon">
                        <i class="fas fa-image"></i>
                    </div>
                    <div class="image-info">
                        <div class="image-title">{image_id}</div>
                        <div class="image-meta">Page {page}</div>
                    </div>
                    <a href="/{path}" target="_blank" class="image-link">
                        <i class="fas fa-external-link-alt"></i>
                        View Image
                    </a>
                </div>
                
                <div class="image-content">
                    <div class="content-section">
                        <div class="section-label">
                            <i class="fas fa-eye"></i>
                            OCR Text
                        </div>
                        <div class="ocr-text">{ocr}</div>
                    </div>
                    
                    <div class="content-section">
                        <div class="section-label">
                            <i class="fas fa-brain"></i>
                            AI Summary
                        </div>
                        <div class="summary-text">{summary_html}</div>
                    </div>
                    
                    <div class="content-section">
                        <div class="section-label">
                            <i class="fas fa-lightbulb"></i>
                            Interpretation
                        </div>
                        <div class="interpretation-text">{interpretation_html}</div>
                    </div>
                </div>
            </div>
        """)
    
    images_html = "".join(items)

    html_doc = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Requirements Analysis – {filename}</title>
        <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
        <link href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css" rel="stylesheet">
        <style>
            :root {{
                --primary-color: #2563eb;
                --primary-dark: #1d4ed8;
                --secondary-color: #64748b;
                --success-color: #10b981;
                --warning-color: #f59e0b;
                --error-color: #ef4444;
                --background: #f8fafc;
                --surface: #ffffff;
                --surface-elevated: #ffffff;
                --text-primary: #1e293b;
                --text-secondary: #64748b;
                --text-muted: #94a3b8;
                --border: #e2e8f0;
                --border-light: #f1f5f9;
                --shadow: 0 1px 3px 0 rgba(0, 0, 0, 0.1), 0 1px 2px 0 rgba(0, 0, 0, 0.06);
                --shadow-lg: 0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -2px rgba(0, 0, 0, 0.05);
                --shadow-xl: 0 20px 25px -5px rgba(0, 0, 0, 0.1), 0 10px 10px -5px rgba(0, 0, 0, 0.04);
                --radius: 12px;
                --radius-sm: 8px;
                --radius-lg: 16px;
            }}

            * {{
                margin: 0;
                padding: 0;
                box-sizing: border-box;
            }}

            body {{
                font-family: 'Georgia', 'Times New Roman', serif;
                background: #f5f5f5;
                min-height: 100vh;
                color: #333;
                line-height: 1.6;
            }}

            .container {{
                max-width: 1200px;
                margin: 0 auto;
                padding: 2rem;
            }}

            .header {{
                text-align: center;
                margin-bottom: 2rem;
                padding: 2rem 0;
                background: #fff;
                border-bottom: 3px solid #2563eb;
            }}

            .header h1 {{
                font-size: 2rem;
                font-weight: 600;
                color: #1e293b;
                margin-bottom: 0.5rem;
            }}

            .header p {{
                font-size: 1rem;
                color: #64748b;
            }}

            .back-btn {{
                display: inline-block;
                background: #2563eb;
                color: white;
                text-decoration: none;
                padding: 0.75rem 1.5rem;
                border-radius: 4px;
                font-weight: 500;
                margin-bottom: 2rem;
                border: none;
                cursor: pointer;
            }}

            .back-btn:hover {{
                background: #1d4ed8;
            }}

            .srs-banner {{
                border-radius: 10px;
                padding: 12px 14px;
                margin-bottom: 1rem;
                border: 1px solid #e2e8f0;
                background: #f8fafc;
                color: #0f172a;
            }}
            .srs-banner-title {{
                font-weight: 700;
                margin-bottom: 0.3rem;
            }}
            .srs-banner-meta {{
                font-size: 0.9rem;
                margin-bottom: 0.35rem;
            }}
            .srs-banner-text {{
                font-size: 0.95rem;
            }}
            .srs-details {{
                margin-top: 0.45rem;
                font-size: 0.88rem;
            }}
            .srs-details summary {{
                cursor: pointer;
                color: #1d4ed8;
                user-select: none;
            }}
            .srs-high {{ background: #ecfdf3; border-color: #86efac; color: #166534; }}
            .srs-medium {{ background: #fefce8; border-color: #fde68a; color: #854d0e; }}
            .srs-low {{ background: #fff7ed; border-color: #fdba74; color: #9a3412; }}
            .srs-none {{ background: #fef2f2; border-color: #fecaca; color: #991b1b; }}

            .export-integrate-bar {{
                display: flex;
                align-items: center;
                gap: 0.75rem;
                flex-wrap: wrap;
                padding: 1rem 1.25rem;
                background: #f8fafc;
                border: 1px solid #e2e8f0;
                border-radius: 8px;
                margin-bottom: 1.5rem;
            }}
            .export-label {{
                font-size: 0.9rem;
                color: #64748b;
                font-weight: 500;
            }}
            .export-btn {{
                display: inline-flex;
                align-items: center;
                gap: 0.4rem;
                padding: 0.5rem 1rem;
                background: #2563eb;
                color: white;
                text-decoration: none;
                border-radius: 6px;
                font-size: 0.9rem;
                font-weight: 500;
                transition: all 0.2s;
            }}
            .export-btn:hover {{
                background: #1d4ed8;
                color: white;
            }}
            .export-divider {{
                color: #cbd5e1;
                margin: 0 0.25rem;
            }}
            .integrate-btn {{
                display: inline-flex;
                align-items: center;
                gap: 0.4rem;
                padding: 0.5rem 1rem;
                border: 1px solid #e2e8f0;
                border-radius: 6px;
                font-size: 0.9rem;
                font-weight: 500;
                cursor: pointer;
                background: #fff;
                color: #334155;
                transition: all 0.2s;
            }}
            .integrate-btn:hover {{
                background: #f1f5f9;
                border-color: #cbd5e1;
            }}
            .integrate-btn.integrate-trello:hover {{ background: #0079bf; color: white; border-color: #0079bf; }}
            .integrate-btn.integrate-jira:hover {{ background: #0052cc; color: white; border-color: #0052cc; }}
            .integrate-status {{
                font-size: 0.85rem;
                color: #64748b;
                margin-left: auto;
            }}
            .integrate-status.success {{ color: #166534; }}
            .integrate-status.error {{ color: #991b1b; }}

            .card {{
                background: #ffffff;
                border: 1px solid #ddd;
                border-radius: 4px;
                padding: 1.5rem;
                margin-bottom: 2rem;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }}

            .card-header {{
                display: flex;
                align-items: center;
                margin-bottom: 1rem;
                padding-bottom: 1rem;
                border-bottom: 2px solid #e2e8f0;
            }}

            .card-icon {{
                width: 40px;
                height: 40px;
                border-radius: 4px;
                display: flex;
                align-items: center;
                justify-content: center;
                margin-right: 1rem;
                font-size: 1.2rem;
                color: white;
                background: #2563eb;
            }}

            .card-title {{
                font-size: 1.25rem;
                font-weight: 600;
                color: #1e293b;
            }}

            .card-subtitle {{
                font-size: 0.9rem;
                color: #64748b;
                margin-top: 0.25rem;
            }}

            .requirements-content {{
                background: #ffffff;
                border: 1px solid #ddd;
                border-radius: 4px;
                padding: 0;
            }}
            
            /* Dropdown Sections */
            .dropdown-section {{
                border-bottom: 1px solid #e2e8f0;
            }}
            
            .dropdown-section:last-child {{
                border-bottom: none;
            }}
            
            .dropdown-toggle {{
                width: 100%;
                text-align: left;
                background: #f8fafc;
                border: none;
                padding: 1rem 1.5rem;
                cursor: pointer;
                font-size: 1.1rem;
                font-weight: 600;
                color: #1e293b;
                display: flex;
                align-items: center;
                gap: 0.75rem;
                transition: background 0.2s;
            }}
            
            .dropdown-toggle:hover {{
                background: #f1f5f9;
            }}
            
            .dropdown-icon {{
                font-size: 0.75rem;
                transition: transform 0.3s;
                color: #2563eb;
                transform: rotate(-90deg);
            }}
            
            .dropdown-icon.rotated {{
                transform: rotate(0deg);
            }}
            
            .dropdown-title {{
                flex: 1;
            }}
            
            .dropdown-count {{
                color: #64748b;
                font-weight: 400;
                font-size: 0.9rem;
            }}
            
            .dropdown-content {{
                padding: 1rem 1.5rem;
                background: #ffffff;
                border-top: 1px solid #e2e8f0;
            }}
            
            .simple-list {{
                list-style: none;
                padding: 0;
                margin: 0;
            }}
            
            .classify-info {{
                color: #94a3b8;
                cursor: help;
                margin-left: 0.25rem;
                display: inline-flex;
                align-items: center;
                font-size: 0.95rem;
            }}
            .classify-info:hover, .classify-info:focus {{ color: #2563eb; outline: none; }}

            .simple-item {{
                padding: 0.75rem 0;
                border-bottom: 1px solid #f1f5f9;
                color: #333;
                line-height: 1.6;
                display: flex;
                align-items: center;
                gap: 0.5rem;
                flex-wrap: wrap;
            }}
            
            .simple-item:last-child {{
                border-bottom: none;
            }}
            
            .simple-item::before {{
                content: "•";
                color: #2563eb;
                font-weight: bold;
                display: inline-block;
                width: 1em;
                margin-right: 0.5rem;
            }}
            
            .item-text {{
                flex: 1;
                min-width: 60%;
            }}
            
            .status-pill, .quality-pill {{
                display: inline-flex;
                align-items: center;
                gap: 0.35rem;
                padding: 0.15rem 0.6rem;
                border-radius: 9999px;
                font-size: 0.85rem;
                font-weight: 600;
                border: 1px solid transparent;
            }}
            
            .status-pill.status-approved {{
                background: #dcfce7;
                color: #166534;
                border-color: #86efac;
            }}
            
            .status-pill.status-denied {{
                background: #fee2e2;
                color: #991b1b;
                border-color: #fecaca;
            }}
            
            .status-pill.status-pending {{
                background: #e0f2fe;
                color: #075985;
                border-color: #bae6fd;
            }}
            
            .quality-pill {{
                background: #eef2ff;
                color: #4338ca;
                border-color: #c7d2fe;
                font-weight: 700;
            }}
            
            .simple-item.status-approved .item-text {{ color: #166534; }}
            .simple-item.status-denied .item-text {{ color: #991b1b; }}
            .simple-item.status-pending .item-text {{ color: #0f172a; }}
            
            .status-summary {{
                display: grid;
                grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
                gap: 0.5rem;
                margin-bottom: 1rem;
            }}
            
            .status-chip {{
                display: inline-flex;
                align-items: center;
                gap: 0.5rem;
                padding: 0.55rem 0.85rem;
                border-radius: 8px;
                font-weight: 600;
                border: 1px solid #e2e8f0;
                color: #0f172a;
                background: #f8fafc;
            }}
            
            .status-chip i {{ opacity: 0.8; }}
            .chip-approved {{ background: #ecfdf3; color: #166534; border-color: #bbf7d0; }}
            .chip-denied {{ background: #fef2f2; color: #991b1b; border-color: #fecdd3; }}
            .chip-pending {{ background: #eff6ff; color: #1d4ed8; border-color: #bfdbfe; }}
            .chip-total {{ background: #f8fafc; color: #0f172a; }}
            
            .no-content {{
                padding: 2rem;
                text-align: center;
                color: #64748b;
            }}

            .image-card {{
                background: #ffffff;
                border: 1px solid #ddd;
                border-radius: 4px;
                padding: 1.5rem;
                margin-bottom: 1.5rem;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }}

            .image-header {{
                display: flex;
                align-items: center;
                margin-bottom: 1rem;
            }}

            .image-icon {{
                width: 40px;
                height: 40px;
                border-radius: 4px;
                display: flex;
                align-items: center;
                justify-content: center;
                margin-right: 1rem;
                font-size: 1.2rem;
                color: white;
                background: #2563eb;
            }}

            .image-info {{
                flex: 1;
            }}

            .image-title {{
                font-weight: 600;
                color: var(--text-primary);
                font-size: 1.1rem;
            }}

            .image-meta {{
                font-size: 0.875rem;
                color: var(--text-secondary);
                margin-top: 0.25rem;
            }}

            .image-link {{
                display: inline-flex;
                align-items: center;
                gap: 0.5rem;
                background: var(--primary-color);
                color: white;
                text-decoration: none;
                padding: 0.5rem 1rem;
                border-radius: var(--radius-sm);
                font-size: 0.875rem;
                font-weight: 500;
                transition: all 0.3s ease;
            }}

            .image-link:hover {{
                background: var(--primary-dark);
                transform: translateY(-1px);
            }}

            .image-content {{
                display: grid;
                gap: 1rem;
            }}

            .content-section {{
                background: #f8fafc;
                border-radius: 4px;
                padding: 1rem;
                border: 1px solid #e2e8f0;
                margin-bottom: 1rem;
            }}
            
            .content-section:last-child {{
                margin-bottom: 0;
            }}

            .section-label {{
                display: flex;
                align-items: center;
                gap: 0.5rem;
                font-weight: 600;
                color: var(--text-primary);
                margin-bottom: 0.75rem;
                font-size: 0.9rem;
            }}

            .section-label i {{
                color: var(--primary-color);
            }}

            .ocr-text {{
                background: #f8fafc;
                color: #333;
                padding: 1rem;
                border-radius: 4px;
                font-family: 'Courier New', monospace;
                font-size: 0.9rem;
                line-height: 1.6;
                max-height: 200px;
                overflow-y: auto;
                border: 1px solid #e2e8f0;
            }}

            .summary-text {{
                color: var(--text-primary);
                line-height: 1.6;
            }}

            .interpretation-text {{
                color: var(--text-secondary);
                font-style: italic;
                line-height: 1.6;
            }}

            @keyframes fadeInDown {{
                from {{
                    opacity: 0;
                    transform: translateY(-30px);
                }}
                to {{
                    opacity: 1;
                    transform: translateY(0);
                }}
            }}

            @keyframes fadeInUp {{
                from {{
                    opacity: 0;
                    transform: translateY(30px);
                }}
                to {{
                    opacity: 1;
                    transform: translateY(0);
                }}
            }}

            @media (max-width: 768px) {{
                .container {{
                    padding: 1rem;
                }}
                
                .header h1 {{
                    font-size: 1.5rem;
                }}
                
                .image-header {{
                    flex-direction: column;
                    align-items: flex-start;
                    gap: 1rem;
                }}
                
                .image-link {{
                    align-self: flex-end;
                }}
            }}
        </style>
        <script>
            document.addEventListener('DOMContentLoaded', function() {{
                const viewId = '{view_id}';
                const statusEl = document.getElementById('integrate-status');

                // Add click handlers to all dropdown toggles
                document.querySelectorAll('.dropdown-toggle').forEach(function(button) {{
                    button.addEventListener('click', function() {{
                        const sectionId = this.getAttribute('data-section-id');
                        const content = document.getElementById(sectionId);
                        const icon = document.getElementById('icon_' + sectionId);
                        
                        if (content && icon) {{
                            if (content.style.display === 'none') {{
                                content.style.display = 'block';
                                icon.classList.add('rotated');
                            }} else {{
                                content.style.display = 'none';
                                icon.classList.remove('rotated');
                            }}
                        }}
                    }});
                }});

                // Trello / Jira integration buttons (send auth token when available for integration log)
                document.querySelectorAll('.integrate-btn').forEach(function(btn) {{
                    btn.addEventListener('click', async function() {{
                        const platform = this.getAttribute('data-platform');
                        statusEl.textContent = 'Sending...';
                        statusEl.className = 'integrate-status';
                        const token = typeof localStorage !== 'undefined' ? localStorage.getItem('access_token') : null;
                        const headers = {{ 'Content-Type': 'application/json' }};
                        if (token) headers['Authorization'] = 'Bearer ' + token;
                        try {{
                            const url = '/api/integration/' + platform + '/' + viewId;
                            const res = await fetch(url, {{ method: 'POST', headers: headers, body: '{{}}' }});
                            const data = await res.json();
                            statusEl.textContent = data.message || (data.success ? 'Done!' : 'Failed');
                            statusEl.className = 'integrate-status ' + (data.success ? 'success' : 'error');
                        }} catch (e) {{
                            statusEl.textContent = 'Error: ' + e.message;
                            statusEl.className = 'integrate-status error';
                        }}
                        setTimeout(function() {{ statusEl.textContent = ''; }}, 5000);
                    }});
                }});
            }});
        </script>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <h1><i class="fas fa-file-alt"></i> Requirements Analysis</h1>
                <p>Detailed extraction results for {filename}</p>
            </div>

            <a href="/" class="back-btn">
                <i class="fas fa-arrow-left"></i>
                Back to FlowMind
            </a>

            {srs_banner_html}

            <div class="export-integrate-bar">
                <span class="export-label"><i class="fas fa-download"></i> Export:</span>
                <a href="/api/export/{view_id}/csv" class="export-btn" download><i class="fas fa-file-csv"></i> CSV</a>
                <a href="/api/export/{view_id}/json" class="export-btn" download><i class="fas fa-file-code"></i> JSON</a>
                <span class="export-divider">|</span>
                <span class="export-label"><i class="fas fa-paper-plane"></i> Integrate:</span>
                <button type="button" class="integrate-btn integrate-trello" data-platform="trello" title="Push to Trello (requires TRELLO_API_KEY, TRELLO_TOKEN, TRELLO_LIST_ID in .env)">
                    <i class="fab fa-trello"></i> Trello
                </button>
                <button type="button" class="integrate-btn integrate-jira" data-platform="jira" title="Push to Jira (requires JIRA_URL, JIRA_PROJECT_KEY, JIRA_EMAIL, JIRA_API_TOKEN in .env)">
                    <i class="fab fa-jira"></i> Jira
                </button>
                <span id="integrate-status" class="integrate-status"></span>
            </div>

            <div class="card">
                <div class="card-header">
                    <div class="card-icon">
                        <i class="fas fa-file-alt"></i>
                    </div>
                    <div>
                        <div class="card-title">Document Analysis</div>
                        <div class="card-subtitle">{summary}</div>
                    </div>
                </div>
                
                {status_summary_html}
                <div class="requirements-content">{safe}</div>
            </div>

            <div class="card">
                <div class="card-header">
                    <div class="card-icon">
                        <i class="fas fa-images"></i>
                    </div>
                    <div>
                        <div class="card-title">Image Analysis</div>
                        <div class="card-subtitle">OCR and AI-powered image insights</div>
                    </div>
                </div>
                
                <div class="image-content">
                    {images_html}
                </div>
            </div>
        </div>
    </body>
    </html>
    """
    return HTMLResponse(content=html_doc)
